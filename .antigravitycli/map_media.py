from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

pptx_path = r"C:\Users\a\Desktop\쇼핑몰 프로젝트 자료모음\6_자소서, 포폴 1차 완성본_260525\포트폴리오(김민도).pptx"
prs = Presentation(pptx_path)

output_file = r"C:\Users\a\Desktop\coding\forGithub\kmd192kmd.github.io\.antigravitycli\pptx_image_mapping.txt"

with open(output_file, 'w', encoding='utf-8') as f:
    for idx, slide in enumerate(prs.slides):
        f.write(f"=== Slide {idx + 1} ===\n")
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text
        f.write(f"  Title: {title}\n")
        for s_idx, shape in enumerate(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Get image name
                image = shape.image
                f.write(f"  Shape {s_idx} [PICTURE]: size={image.size}, content_type={image.content_type}\n")
        f.write("\n")

print("Image mapping output written successfully!")
