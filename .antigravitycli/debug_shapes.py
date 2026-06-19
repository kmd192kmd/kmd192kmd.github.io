from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

pptx_path = r"C:\Users\a\Desktop\쇼핑몰 프로젝트 자료모음\6_자소서, 포폴 1차 완성본_260525\포트폴리오(김민도).pptx"
prs = Presentation(pptx_path)

for idx, slide in enumerate(prs.slides):
    print(f"=== Slide {idx + 1} ===")
    for s_idx, shape in enumerate(slide.shapes):
        print(f"  Shape {s_idx}: name='{shape.name}', type={shape.shape_type}")
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            print(f"    Group contains {len(shape.shapes)} shapes")
            for sub_shape in shape.shapes:
                print(f"      Sub-shape: name='{sub_shape.name}', type={sub_shape.shape_type}")
        if shape.fill and hasattr(shape.fill, 'type'):
            print(f"    Fill type: {shape.fill.type}")
    if idx >= 5: # Just first few slides for debugging
        break
