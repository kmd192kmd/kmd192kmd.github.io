"""
포트폴리오(박지명)_v1.pptx - 박지명 국비 지원자 기반 디자인 v1
Theme: "Clean Minimalist Light Editorial Design"
- 배경: 깨끗한 화이트 & 라이트 그레이
- 목차: 세로형 2열 분할 그리드 카드 레이아웃
- 개요: 좌측 스펙 카드 + 우측 스크린샷 & 기획의도
- 핵심 구현: 3열 2행 카드 그리드 레이아웃
- 트러블슈팅: PROBLEM & CONTEXT / SOLUTION & RESOLUTION 세로 적층형 좌측 카드 + 우측 이미지 프레임
- 회고: 좌측 다크 테마의 Quote 블록 + 우측 세로 Lessons 카드덱
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── DESIGN TOKENS ───
DARK         = RGBColor(15, 23, 42)      # Deep Slate (slate-900)
BG_LIGHT     = RGBColor(255, 255, 255)   # White
WHITE        = RGBColor(255, 255, 255)   # White
GRAY50       = RGBColor(248, 250, 252)   # slate-50
GRAY100      = RGBColor(241, 245, 249)   # slate-100
GRAY200      = RGBColor(226, 232, 240)   # slate-200
GRAY300      = RGBColor(203, 213, 225)   # slate-300
GRAY400      = RGBColor(148, 163, 184)   # slate-400
GRAY500      = RGBColor(100, 116, 139)   # slate-500
GRAY700      = RGBColor(51, 65, 85)      # slate-700
GRAY900      = RGBColor(15, 23, 42)

# Brand Colors (Park Ji-myung)
P1_MAIN  = RGBColor(99, 102, 241)    # Indigo-500 (제시요)
P1_LINE  = RGBColor(199, 210, 254)   # Indigo-200

P2_MAIN  = RGBColor(16, 185, 129)    # Emerald-500 (커뮤니티)
P2_LINE  = RGBColor(167, 243, 208)   # Emerald-200

FN = "Malgun Gothic"

# Image dirs
EXT = r"C:\Users\a\Desktop\coding\forGithub\kmd192kmd.github.io\images\extracted"

def ip(name):
    p = os.path.join(EXT, name)
    return p if os.path.exists(p) else None

# ─── PRIMITIVES ───
def bg(slide, c):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = c

def rect(slide, l, t, w, h, fill=None, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
    else:
        s.line.fill.background()
    return s

def rrect(slide, l, t, w, h, fill=None, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
    else:
        s.line.fill.background()
    return s

def tx(slide, l, t, w, h):
    return slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))

def rn(p, text, size=12, bold=False, color=GRAY900, italic=False):
    r = p.add_run(); r.text = text; r.font.name = FN
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.italic = italic
    return r

def pic(slide, path, l, t, w, h, cap=None):
    if path and os.path.exists(path):
        try:
            slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
            if cap:
                tb = tx(slide, l, t + h + 0.05, w, 0.3)
                p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
                rn(p, cap, 9, italic=True, color=GRAY500)
        except Exception as e:
            print(f"Picture error: {e}")

# ─── TECH PILLS DRAW FOR DIVIDERS ───
def draw_tech_pills(slide, start_x, start_y, pills_list, fill_color, text_color):
    curr_x = start_x
    curr_y = start_y
    for pill in pills_list:
        pill_w = len(pill) * 0.085 + 0.4
        if curr_x + pill_w > 12.5:
            curr_x = start_x
            curr_y += 0.45
        
        p_shp = rrect(slide, curr_x, curr_y, pill_w, 0.35, DARK, text_color)
        tf = p_shp.text_frame
        tf.word_wrap = False
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.03)
        tf.margin_bottom = Inches(0.03)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        rn(p, pill, 9.5, True, text_color)
        curr_x += pill_w + 0.12

# ─── MINIMAL CLEAN HEADER ───
PROJ_NAMES = ["Jesiyo (제시요)", "Community (커뮤니티)"]
PROJ_COLORS = [P1_MAIN, P2_MAIN]

def add_clean_header(slide, active_proj_idx, section_name, title, accent):
    rect(slide, 0.8, 0.45, 0.08, 0.65, accent)
    tb = tx(slide, 1.05, 0.35, 8.5, 0.85)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    rn(p, f"{section_name.upper()}   //\n", 9.5, True, accent)
    rn(p, title, 20, True, GRAY900)
    
    tb_r = tx(slide, 9.5, 0.45, 3.033, 0.65)
    tf_r = tb_r.text_frame
    tf_r.margin_left = tf_r.margin_top = tf_r.margin_right = tf_r.margin_bottom = 0
    p_r = tf_r.paragraphs[0]; p_r.alignment = PP_ALIGN.RIGHT
    rn(p_r, PROJ_NAMES[active_proj_idx].upper(), 11, True, accent)
    rn(p_r, "  |  ", 10, False, GRAY300)
    rn(p_r, section_name.split()[-1], 10, True, GRAY400)

# ─── PARSE BULLET WITH INLINE MARKERS ───
def parse_bullet_text(tf, items, accent, is_first_p=True):
    for item in items:
        p = tf.paragraphs[0] if is_first_p else tf.add_paragraph()
        is_first_p = False
        p.space_after = Pt(7); p.space_before = Pt(3)
        rn(p, "▪ ", 11, True, accent)
        
        if "**" in item:
            parts = item.split("**")
            for pi, part in enumerate(parts):
                if pi % 2 == 1:
                    rn(p, part, 11, True, GRAY900)
                else:
                    rn(p, part, 11, color=GRAY700)
        else:
            rn(p, item, 11, color=GRAY700)

# ─── FLOATING FOOTER ───
def add_footer(slide, active_page, accent):
    rect(slide, 0.8, 6.95, 11.733, 0.015, GRAY200)
    lt = tx(slide, 0.8, 7.02, 6.0, 0.3)
    lp = lt.text_frame.paragraphs[0]
    rn(lp, "PARK JI MYUNG  ·  BACKEND DEVELOPER PORTFOLIO", 8, True, GRAY400)
    
    pt = tx(slide, 9.0, 7.02, 3.533, 0.3)
    pp = pt.text_frame.paragraphs[0]; pp.alignment = PP_ALIGN.RIGHT
    rn(pp, active_page, 9.5, True, accent)

# ─── SECTION DIVIDER (간지) SLIDE ───
def divider(part_num, title, subtitle, accent, accent_lite, pills_list=None):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, DARK)
    
    rect(sl, 4.6, 0, 0.08, 7.5, accent)
    
    pt = tx(sl, 0.6, 1.8, 3.5, 0.4)
    pp = pt.text_frame.paragraphs[0]
    rn(pp, f"PART 0{part_num}.", 12, True, accent)
    
    nt = tx(sl, 0.6, 2.3, 3.5, 2.5)
    np = nt.text_frame.paragraphs[0]
    rn(np, f"0{part_num}", 100, True, WHITE)
    
    cat_lbls = ["REALTIME BID & CHAT", "JSP MVC COMMUNITY"]
    ct = tx(sl, 0.6, 5.0, 3.5, 0.4)
    cp = ct.text_frame.paragraphs[0]
    rn(cp, cat_lbls[part_num - 1], 10.5, True, GRAY400)
    
    tt = tx(sl, 5.2, 2.0, 7.5, 1.0)
    tp = tt.text_frame.paragraphs[0]
    rn(tp, title, 36, True, WHITE)
    
    st = tx(sl, 5.2, 3.0, 7.2, 1.0)
    st.text_frame.word_wrap = True
    sp = st.text_frame.paragraphs[0]
    rn(sp, subtitle, 14, color=GRAY300)
    
    rect(sl, 5.2, 4.3, 7.2, 0.015, RGBColor(38, 54, 84))
    
    if pills_list:
        draw_tech_pills(sl, 5.2, 4.7, pills_list, None, accent)

# ─── OVERVIEW LAYOUT ───
def overview_slide(project_idx, section, title, items, accent, img_path=None, img_caption=None, active_page=""):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, BG_LIGHT)
    add_clean_header(sl, project_idx, section, title, accent)
    add_footer(sl, active_page, accent)
    
    specs = []
    desc_items = []
    for item in items:
        if any(x in item for x in ["**기간**", "**인력구성**", "**역할**", "**기술 스택**", "**기획 의도**"]):
            specs.append(item)
        else:
            desc_items.append(item)
            
    rrect(sl, 0.8, 1.6, 3.8, 5.0, GRAY50, GRAY200)
    rect(sl, 0.8, 1.6, 3.8, 0.08, accent)
    
    tb_spec = tx(sl, 1.05, 1.85, 3.3, 4.5)
    tf_s = tb_spec.text_frame; tf_s.word_wrap = True
    tf_s.margin_left = tf_s.margin_top = tf_s.margin_right = tf_s.margin_bottom = 0
    
    p_h = tf_s.paragraphs[0]; p_h.space_after = Pt(12)
    rn(p_h, "PROJECT SPEC", 13, True, GRAY900)
    
    is_first = True
    for spec in specs:
        p = tf_s.paragraphs[0] if is_first else tf_s.add_paragraph()
        is_first = False
        p.space_after = Pt(8); p.space_before = Pt(2)
        
        if "  " in spec:
            key, val = spec.split("  ", 1)
        elif "** " in spec:
            key, val = spec.split("** ", 1)
            key = key + "**"
        else:
            key, val = spec, ""
            
        k_clean = key.replace("**", "").strip()
        rn(p, f"▪ {k_clean}\n", 10, True, accent)
        rn(p, val.strip(), 9.5, color=GRAY700)
        
    has_img = img_path and os.path.exists(img_path)
    if has_img:
        rrect(sl, 4.9, 1.6, 7.633, 3.2, GRAY50, GRAY200)
        pic(sl, img_path, 5.05, 1.75, 7.333, 2.5)
        if img_caption:
            tb_cap = tx(sl, 5.05, 4.3, 7.333, 0.3)
            p = tb_cap.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            rn(p, img_caption, 8.5, italic=True, color=GRAY500)
            
        tb_desc = tx(sl, 4.9, 4.9, 7.633, 1.7)
        tf_d = tb_desc.text_frame; tf_d.word_wrap = True
        tf_d.margin_left = tf_d.margin_top = tf_d.margin_bottom = tf_d.margin_right = 0
        parse_bullet_text(tf_d, desc_items, accent, is_first_p=True)
    else:
        tb_desc = tx(sl, 4.9, 1.6, 7.633, 5.0)
        tf_d = tb_desc.text_frame; tf_d.word_wrap = True
        tf_d.margin_left = tf_d.margin_top = tf_d.margin_bottom = tf_d.margin_right = 0
        parse_bullet_text(tf_d, desc_items, accent, is_first_p=True)

# ─── CORE FEATURES GRID LAYOUT ───
def features_slide(project_idx, section, title, items, accent, active_page=""):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, BG_LIGHT)
    add_clean_header(sl, project_idx, section, title, accent)
    add_footer(sl, active_page, accent)
    
    n = len(items)
    cols = 3 if n <= 6 else 4
    rows = 2
        
    for idx, item in enumerate(items):
        r = idx // cols
        c = idx % cols
        
        if cols == 3:
            w_card = 3.68; gap = 0.35; l_start = 0.8
        else:
            w_card = 2.68; gap = 0.32; l_start = 0.8
            
        x = l_start + c * (w_card + gap)
        y = 1.6 + r * 2.55
        h_card = 2.3
        
        rrect(sl, x, y, w_card, h_card, GRAY50, GRAY200)
        rect(sl, x, y, w_card, 0.06, accent)
        
        tb = tx(sl, x + 0.15, y + 0.15, w_card - 0.3, h_card - 0.3)
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        text = item
        if text[0].isdigit() and text[1:3] in [". ", ") ", "  "]:
            text = text[3:]
        elif text[0].isdigit() and text[1] == " ":
            text = text[2:]
            
        card_num = f"0{idx + 1}"
        p = tf.paragraphs[0]; p.space_after = Pt(4)
        rn(p, f"{card_num}  ", 10.5, True, accent)
        
        if "**" in text:
            parts = text.split("**")
            title_text = parts[1]
            desc_text = "".join(parts[2:]).strip(" —").strip()
            rn(p, title_text, 11, True, GRAY900)
            p2 = tf.add_paragraph(); p2.space_before = Pt(2)
            rn(p2, desc_text, 9.5, color=GRAY700)
        else:
            rn(p, text, 10.5, color=GRAY700)

# ─── TECHNICAL DETAIL LAYOUT ───
def technical_slide(project_idx, section, title, items, accent, img_path=None, img_caption=None, active_page=""):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, BG_LIGHT)
    add_clean_header(sl, project_idx, section, title, accent)
    add_footer(sl, active_page, accent)
    
    has_img = img_path and os.path.exists(img_path)
    if has_img:
        tb_text = tx(sl, 0.8, 1.6, 5.8, 5.0)
        tf = tb_text.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        parse_bullet_text(tf, items, accent, is_first_p=True)
        
        rrect(sl, 6.9, 1.6, 5.633, 5.0, GRAY50, GRAY200)
        pic(sl, img_path, 7.05, 1.75, 5.333, 4.0)
        if img_caption:
            tb_cap = tx(sl, 7.05, 5.85, 5.333, 0.35)
            p = tb_cap.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            rn(p, img_caption, 9, italic=True, color=GRAY500)
    else:
        if len(items) >= 5:
            mid = (len(items) + 1) // 2
            left_items = items[:mid]
            right_items = items[mid:]
            
            ct_l = tx(sl, 0.8, 1.6, 5.5, 5.0)
            tf_l = ct_l.text_frame; tf_l.word_wrap = True
            tf_l.margin_left = tf_l.margin_top = tf_l.margin_bottom = tf_l.margin_right = 0
            parse_bullet_text(tf_l, left_items, accent, is_first_p=True)
            
            rect(sl, 6.666, 1.8, 0.015, 4.5, GRAY200)
            
            ct_r = tx(sl, 7.0, 1.6, 5.5, 5.0)
            tf_r = ct_r.text_frame; tf_r.word_wrap = True
            tf_r.margin_left = tf_r.margin_top = tf_r.margin_bottom = tf_r.margin_right = 0
            parse_bullet_text(tf_r, right_items, accent, is_first_p=True)
        else:
            ct = tx(sl, 0.8, 1.6, 11.733, 5.0)
            tf = ct.text_frame; tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
            parse_bullet_text(tf, items, accent, is_first_p=True)

# ─── TROUBLESHOOTING LAYOUT ───
def split_trouble_items(items):
    for idx in range(1, len(items)):
        if any(kw in items[idx] for kw in ["**해결**", "**개선**", "주입", "통일", "적용"]):
            return items[:idx], items[idx:]
    mid = (len(items) + 1) // 2
    return items[:mid], items[mid:]

def troubleshooting_slide(project_idx, section, title, items, accent, img_path=None, img_caption=None, active_page=""):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, BG_LIGHT)
    add_clean_header(sl, project_idx, section, title, accent)
    add_footer(sl, active_page, accent)
    
    left_items, right_items = split_trouble_items(items)
    has_img = img_path and os.path.exists(img_path)
    
    if has_img:
        rrect(sl, 0.8, 1.6, 5.8, 2.3, GRAY50, GRAY200)
        rect(sl, 0.8, 1.6, 0.08, 2.3, accent)
        tb_p = tx(sl, 1.05, 1.75, 5.35, 2.05)
        tf_p = tb_p.text_frame; tf_p.word_wrap = True
        tf_p.margin_left = tf_p.margin_top = tf_p.margin_bottom = tf_p.margin_right = 0
        p_ph = tf_p.paragraphs[0]; p_ph.space_after = Pt(4)
        rn(p_ph, "PROBLEM & CONTEXT", 11, True, GRAY900)
        parse_bullet_text(tf_p, left_items, accent, is_first_p=False)
        
        rrect(sl, 0.8, 4.1, 5.8, 2.5, GRAY50, GRAY200)
        rect(sl, 0.8, 4.1, 0.08, 2.5, accent)
        tb_s = tx(sl, 1.05, 4.25, 5.35, 2.25)
        tf_s = tb_s.text_frame; tf_s.word_wrap = True
        tf_s.margin_left = tf_s.margin_top = tf_s.margin_bottom = tf_s.margin_right = 0
        p_sh = tf_s.paragraphs[0]; p_sh.space_after = Pt(4)
        rn(p_sh, "SOLUTION & RESOLUTION", 11, True, GRAY900)
        parse_bullet_text(tf_s, right_items, accent, is_first_p=False)
        
        rrect(sl, 6.9, 1.6, 5.6, 5.0, GRAY50, GRAY200)
        pic(sl, img_path, 7.05, 1.75, 5.3, 4.0)
        if img_caption:
            tb_cap = tx(sl, 7.05, 5.85, 5.3, 0.35)
            p = tb_cap.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            rn(p, img_caption, 8.5, italic=True, color=GRAY500)
    else:
        rrect(sl, 0.8, 1.6, 5.6, 5.0, GRAY50, GRAY200)
        rect(sl, 0.8, 1.6, 0.08, 5.0, accent)
        tb_l = tx(sl, 1.05, 1.8, 5.15, 4.6)
        tf_l = tb_l.text_frame; tf_l.word_wrap = True
        tf_l.margin_left = tf_l.margin_top = tf_l.margin_bottom = tf_l.margin_right = 0
        p_lh = tf_l.paragraphs[0]; p_lh.space_after = Pt(10)
        rn(p_lh, "PROBLEM & CONTEXT", 12, True, GRAY900)
        parse_bullet_text(tf_l, left_items, accent, is_first_p=False)
        
        rrect(sl, 6.9, 1.6, 5.6, 5.0, GRAY50, GRAY200)
        rect(sl, 6.9, 1.6, 0.08, 5.0, accent)
        tb_r = tx(sl, 7.15, 1.8, 5.15, 4.6)
        tf_r = tb_r.text_frame; tf_r.word_wrap = True
        tf_r.margin_left = tf_r.margin_top = tf_r.margin_bottom = tf_r.margin_right = 0
        p_rh = tf_r.paragraphs[0]; p_rh.space_after = Pt(10)
        rn(p_rh, "SOLUTION & RESOLUTION", 12, True, GRAY900)
        parse_bullet_text(tf_r, right_items, accent, is_first_p=False)

# ─── RETROSPECTIVE LAYOUT ───
RETRO_QUOTES = {
    0: "도메인별 설정 분리와 캐싱 처리로\n안정적이고 확장 가능한 실시간 시스템 구축의\n중요성을 실감했습니다.",
    1: "단일 템플릿과 라우팅 최적화로\nDRY 원칙을 지키고, 계층화 아키텍처의\n유지보수 효율을 깊이 학습했습니다."
}

def retrospective_slide(project_idx, section, title, items, accent, active_page=""):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, BG_LIGHT)
    add_clean_header(sl, project_idx, section, title, accent)
    add_footer(sl, active_page, accent)
    
    rrect(sl, 0.8, 1.6, 4.8, 5.0, DARK, None)
    rect(sl, 0.8, 1.6, 0.12, 5.0, accent)
    
    tb_q = tx(sl, 1.2, 2.2, 4.0, 4.0)
    tf_q = tb_q.text_frame; tf_q.word_wrap = True
    tf_q.margin_left = tf_q.margin_top = tf_q.margin_bottom = tf_q.margin_right = 0
    p_q1 = tf_q.paragraphs[0]; p_q1.space_after = Pt(8)
    rn(p_q1, "DEVELOPER'S NOTE", 11, True, accent)
    
    p_q2 = tf_q.add_paragraph(); p_q2.space_before = Pt(8); p_q2.space_after = Pt(12)
    quote_text = RETRO_QUOTES.get(project_idx)
    rn(p_q2, f'"{quote_text}"', 14, False, WHITE, italic=True)
    
    lessons = []
    curr_title = ""
    curr_desc = []
    for item in items:
        if item.startswith("**") or (len(item) < 30 and not item.startswith("-")):
            if curr_title:
                lessons.append((curr_title, " ".join(curr_desc)))
            curr_title = item.replace("**", "").strip()
            curr_desc = []
        else:
            curr_desc.append(item)
    if curr_title:
        lessons.append((curr_title, " ".join(curr_desc)))
        
    num_lessons = len(lessons)
    for idx, (ltitle, ldesc) in enumerate(lessons[:3]):
        if num_lessons >= 3:
            h_card = 1.5; y = 1.6 + idx * 1.73; text_size = 9.5; desc_space = Pt(2)
        elif num_lessons == 2:
            h_card = 2.38; y = 1.6 + idx * 2.62; text_size = 10; desc_space = Pt(4)
        else:
            h_card = 5.0; y = 1.6; text_size = 10.5; desc_space = Pt(6)
            
        rrect(sl, 6.0, y, 6.5, h_card, GRAY50, GRAY200)
        rect(sl, 6.0, y, 6.5, 0.06, accent)
        
        tb_l = tx(sl, 6.2, y + 0.15, 6.1, h_card - 0.25)
        tf_l = tb_l.text_frame; tf_l.word_wrap = True
        tf_l.margin_left = tf_l.margin_top = tf_l.margin_bottom = tf_l.margin_right = 0
        
        p = tf_l.paragraphs[0]; p.space_after = desc_space
        rn(p, f"0{idx + 1}  {ltitle}", 11.5, True, GRAY900)
        p2 = tf_l.add_paragraph()
        rn(p2, ldesc, text_size, color=GRAY700)


# ══════════════════════════════════════════════════════════════════
# SLIDE 1: COVER
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, DARK)

for gx in [1.5, 3.0, 4.5, 6.0, 7.5, 9.0, 10.5, 12.0]:
    rect(sl, gx, 0, 0.01, 7.5, RGBColor(22, 32, 54), None)
for gy in [1.5, 3.0, 4.5, 6.0]:
    rect(sl, 0, gy, 13.333, 0.01, RGBColor(22, 32, 54), None)

rect(sl, 10.1, 0, 0.8, 7.5, P1_MAIN)
rect(sl, 11.0, 0, 0.8, 7.5, P2_MAIN)
rect(sl, 0.8, 1.8, 0.06, 3.8, P1_MAIN)

mt = tx(sl, 1.1, 1.7, 7.8, 4.5)
tf = mt.text_frame; tf.word_wrap = True
p1 = tf.paragraphs[0]; p1.space_after = Pt(12)
rn(p1, "BACKEND ARCHITECTURE & SYSTEM DESIGN", 13, True, P1_MAIN)
p2 = tf.add_paragraph(); p2.space_after = Pt(4)
rn(p2, "PARK JI MYUNG", 52, True, WHITE)
p3 = tf.add_paragraph(); p3.space_after = Pt(24)
rn(p3, "박지명 백엔드 개발자 포트폴리오", 22, color=GRAY300)
p4 = tf.add_paragraph(); p4.space_after = Pt(14)
rn(p4, '"성장을 멈추지 않는 책임감 있는 백엔드 엔지니어"', 14, color=GRAY400, italic=True)
p5 = tf.add_paragraph(); p5.space_after = Pt(30)
rn(p5, "Java Full-Stack 개발 과정을 거쳐 실시간 분산 시스템 설계와 효율적인 데이터 모델링을 주력으로 삼고 있습니다.", 11, color=GRAY300)
p6 = tf.add_paragraph()
rn(p6, "010-7571-0454  ·  wlaudaoq123@naver.com  ·  github.com/wauder0430", 11.5, color=GRAY400)


# ══════════════════════════════════════════════════════════════════
# SLIDE 2: PROFILE & SKILLS
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, BG_LIGHT)

add_clean_header(sl, 0, "ABOUT ME", "프로필 및 보유 기술 스택 요약", P1_MAIN)

# Profile Card
rrect(sl, 0.8, 1.6, 5.6, 5.0, GRAY50, GRAY200)
rect(sl, 0.8, 1.6, 5.6, 0.08, P1_MAIN)
lc_t = tx(sl, 1.1, 1.85, 5.0, 4.5)
ltf = lc_t.text_frame; ltf.word_wrap = True
p_h = ltf.paragraphs[0]; p_h.space_after = Pt(14)
rn(p_h, "Profile", 18, True, GRAY900)

for label, val in [
    ("이름 / 연령", "박지명  /  1996.04.30"),
    ("학력", "전주대학교 컴퓨터공학과 졸업 (학점 4.12 / 4.5)"),
    ("훈련 이수", "SIST 쌍용교육센터 Java Full-Stack 양성과정 (130일, 1040시간)"),
    ("자격증", "정보처리기사 (2025.12)  |  1종보통운전면허"),
    ("기타 활동", "2022 오픈소스 컨트리뷰션 (OpenStack)  |  JJ Startup Boot Camp 2등"),
]:
    pl = ltf.add_paragraph(); pl.space_before = Pt(4)
    rn(pl, f"▸ {label}", 11, True, GRAY900)
    pv = ltf.add_paragraph(); pv.space_after = Pt(2)
    rn(pv, val, 10, color=GRAY700)

# Technical Skills
rrect(sl, 6.9, 1.6, 5.6, 5.0, GRAY50, GRAY200)
rect(sl, 6.9, 1.6, 5.6, 0.08, P1_MAIN)
rc_t = tx(sl, 7.2, 1.85, 5.0, 4.5)
rtf = rc_t.text_frame; rtf.word_wrap = True
s_h = rtf.paragraphs[0]; s_h.space_after = Pt(14)
rn(s_h, "Technical Skills", 18, True, GRAY900)

for cat, itms in [
    ("Languages", "Java · JavaScript · HTML5 · CSS3"),
    ("Frameworks & Libs", "Spring Boot · Spring Framework · Spring Security · Spring AOP · MyBatis · Servlet/JSP"),
    ("Database & Cache", "Oracle Database · Redis · Redisson"),
    ("Tools & DevOps", "Docker · GitHub Action · Git/GitHub · Figma · ERD Cloud"),
]:
    pc = rtf.add_paragraph(); pc.space_before = Pt(6)
    rn(pc, f"▸ {cat}", 11, True, GRAY900)
    pi = rtf.add_paragraph(); pi.space_after = Pt(2)
    rn(pi, itms, 10, color=GRAY700)

add_footer(sl, "02 / 16", P1_MAIN)


# ══════════════════════════════════════════════════════════════════
# SLIDE 3: TABLE OF CONTENTS
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, BG_LIGHT)

add_clean_header(sl, 0, "INDEX", "수행 프로젝트 인덱스 요약", P1_MAIN)

projects = [
    ("PART #1", "Jesiyo", "제시요 (경매·중고거래 플랫폼)", "WebSocket · Redis Pub/Sub 기반 실시간 채팅 + 채팅·입찰 통합 UX\nSpring · WebSocket · Redis · Redisson · Spring AOP · MyBatis · Oracle", P1_MAIN),
    ("PART #2", "Community", "커뮤니티 (JSP Model2)", "MVC Model 2 설계 — 동적 라우팅·쿼리 최적화 적용 웹 게시판\nJava · Servlet/JSP · Oracle · JDBC · Tailwind CSS · Git/GitHub", P2_MAIN),
]

for i, (part, name_en, name_ko, desc, clr) in enumerate(projects):
    x = 0.8 + i * 5.9
    rrect(sl, x, 1.6, 5.6, 5.0, GRAY50, GRAY200)
    rect(sl, x, 1.6, 5.6, 0.08, clr)
    
    tb = tx(sl, x + 0.3, 1.85, 5.0, 4.5)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p = tf.paragraphs[0]; p.space_after = Pt(2)
    rn(p, part, 11, True, clr)
    p2 = tf.add_paragraph(); p2.space_before = Pt(4); p2.space_after = Pt(2)
    rn(p2, name_en, 18, True, GRAY900)
    p3 = tf.add_paragraph(); p3.space_after = Pt(12)
    rn(p3, name_ko, 12, True, GRAY500)
    
    rect(sl, x + 0.3, 2.9, 5.0, 0.015, GRAY200)
    
    p4 = tf.add_paragraph(); p4.space_before = Pt(12)
    rn(p4, desc, 10, False, GRAY700)

add_footer(sl, "03 / 16", P1_MAIN)


# ══════════════════════════════════════════════════════════════════
# PART 1: JESIYO
# ══════════════════════════════════════════════════════════════════
js_pills = ["Spring", "WebSocket", "Redis", "Redisson", "Spring AOP", "MyBatis", "Oracle", "Spring Security"]
divider(1, "Jesiyo (제시요)",
        "WebSocket · Redis Pub/Sub 기반 실시간 채팅 + 채팅·입찰 통합 UX 구현\n화면 전환 부담 해소를 위한 단일 흐름 UI 설계 및 캐시 오프로딩 적용",
        P1_MAIN, P1_LINE, js_pills)

overview_slide(0, "01  개요", "개요 · 기획의도", [
    "**프로젝트 개요**  경매·중고거래 채팅 플랫폼 제시요 (jesiyo)",
    "**기획 의도**  채팅과 입찰이 분리되어 페이지를 오가던 기존 UX를 개선하기 위해, 채팅방 내 경매 탭으로 두 기능을 한 화면에서 처리하는 통합 UX를 설계했습니다. 또한 메시징을 WebSocket·Redis Pub/Sub 구조로 다중 서버 확장이 가능하도록 구축했습니다.",
    "**기간**  분석·설계 2026.04.15 ~ 04.22  ·  구현 2026.04.23 ~ 04.29",
    "**인력구성**  4명 (초기 환경 설정 및 채팅 파트 담당)",
    "**역할**  초기 환경 구성, WebSocket 실시간 채팅, Redis Pub/Sub, AOP 추천, 채팅 캐싱",
    "**기술 스택**  Spring · WebSocket · Redis · Redisson · Spring AOP · MyBatis · Oracle · Spring Security",
], P1_MAIN, ip("image31.png"), "실시간 경매 및 채팅 통합 UI", "05 / 16")

features_slide(0, "02  핵심 구현", "핵심 구현 사항 (담당: 채팅 파트)", [
    "1. **프로젝트 초기 환경 구성**  MyBatis·WebSocket·Redis Config 전담 설정, 도메인별 Redis 분리로 충돌 방지",
    "2. **Redis Pub/Sub 실시간 채팅 아키텍처**  직접 브로드캐스트 대신 Publisher -> Redis 채널 -> Subscriber 구조 설계",
    "3. **Spring AOP 키워드 추천**  @AfterReturning 활용, 채팅 내역 분석해 관련 상품 추천 (채팅 단일 책임 유지)",
    "4. **채팅 내역 Redis 캐싱**  DB 대신 Redis에서 채팅 내역을 먼저 읽어 DB 부하 감소, 최근 50건 캐시 유지",
    "5. **fetch API 비동기 채팅 UI**  목록·방 생성·초대코드 입장·참여자 갱신 등 새로고침 없는 화면 구현",
], P1_MAIN, "06 / 16")

technical_slide(0, "03  데이터 구조", "시스템 아키텍처 & 주요 데이터 흐름", [
    "**WebSocket & Redis Pub/Sub 연동**  다중 WAS 구성에서도 Redis 채널을 통해 실시간 메시지가 전 서버에 브로드캐스팅되는 구조 설계",
    "**Cache-Aside 패턴의 채팅 캐싱**  채팅 기록 조회 시 캐시(Redis) 확인 -> Miss 발생 시 DB 조회 및 캐시 갱신 구조 수립",
    "**AOP 상품 추천 분리**  메시지 송신 완료(Aspect 포인트컷) 감지 시 비동기 추천 엔진을 구동해 채팅 비즈니스 로직과 추천 로직 완벽 격리",
], P1_MAIN, ip("image33.png"), "채팅 시스템 아키텍처 설계", "07 / 16")

troubleshooting_slide(0, "04  트러블슈팅", "주입 충돌 & 세션 연동 트러블슈팅", [
    "**다중 서버 시 메시지 미수신 문제**",
    "- 핸들러가 메모리에만 세션을 담아두어 서버 1대에서만 동작함.",
    "- Redis Pub/Sub으로 전환하여, 어느 WAS로 접속한 사용자든 Redis 채널을 통해 동일한 실시간 메시지를 수신 가능하도록 개선.",
    "**@Qualifier + @RequiredArgsConstructor 주입 충돌**",
    "- Lombok의 자동 생성자가 @Qualifier 어노테이션을 파라미터에 누락시켜 Bean 주입 에러 발생.",
    "- Lombok 어노테이션 대신 생성자 주입을 명시적으로 작성하고 파라미터에 @Qualifier를 직접 선언하여 정확한 Bean 바인딩 달성.",
], P1_MAIN, ip("image37.png"), "채팅 내역 조회 트러블슈팅 지표", "08 / 16")

retrospective_slide(0, "05  회고", "프로젝트 후기 · 회고", [
    "**설정의 중요성 체감**",
    "Redis 설정을 구분 없이 두었을 때 Bean 충돌을 겪으며 도메인별 설정 격리(Bean Name 분리)가 협업과 유지보수에 미치는 막대한 영향력을 깨달았습니다.",
    "**AOP로 관심사 분리**",
    "채팅 전송과 상품 추천은 서로 다른 책임입니다. AOP(@AfterReturning)를 이용해 결합도를 낮춘 덕분에 추천 로직의 잦은 변경에도 채팅 본체 코드는 안전하게 단일 책임을 유지할 수 있었습니다.",
    "**캐싱 처리 시 원본 보호 우선**",
    "캐시는 언제든 유실될 수 있습니다. DB를 항상 '단일 진실 공급원(Single Source of Truth)'으로 두고 Redis는 보조 조속 스토어로 쓰는 설계 철학을 실천했습니다.",
], P1_MAIN, "09 / 16")


# ══════════════════════════════════════════════════════════════════
# PART 2: COMMUNITY
# ══════════════════════════════════════════════════════════════════
cm_pills = ["Java", "Servlet/JSP", "Oracle", "JDBC", "Tailwind CSS", "DaisyUI", "GitHub", "Figma", "ERD Cloud"]
divider(2, "Community (커뮤니티)",
        "Servlet과 JSP MVC Model 2 구조로 구현한 스터디·학습 커뮤니티 플랫폼\nDRY 원칙의 동적 게시판 연동, 복합 제약을 활용한 어뷰징 차단 및 권한 검증 구현",
        P2_MAIN, P2_LINE, cm_pills)

overview_slide(1, "01  개요", "개요 · 기획의도", [
    "**프로젝트 개요**  학습 관리 & 스터디 커뮤니티 (JSP Model2)",
    "**기획 의도**  프레임워크 없이 서블릿과 JSP MVC Model 2 구조를 밑바닥부터 직접 설계함으로써 웹 애플리케이션의 HTTP 생명주기와 데이터 흐름을 깊이 있게 이해하고자 했습니다. 커뮤니티 파트를 단독 전담하여 4종 이상의 게시판을 유연하게 제어하는 시스템을 구축했습니다.",
    "**기간**  2026.03.19 ~ 2026.03.31 (약 2주)",
    "**인력구성**  4명 (커뮤니티 도메인 단독 전담)",
    "**역할**  커뮤니티 도메인 전담, 페이징·검색 필터링 로직, 공통 권한 Filter 설계",
    "**기술 스택**  Java · Servlet/JSP · Oracle · JDBC · HTML/CSS/JS · Tailwind/DaisyUI · GitHub · Figma · ERD Cloud",
], P2_MAIN, ip("image20.png"), "스터디 목록 화면", "11 / 16")

features_slide(1, "02  핵심 구현", "핵심 구현 사항 (담당: 커뮤니티)", [
    "1. **DRY 원칙 기반 동적 다중 게시판**  4종 이상의 게시판을 단일 JSP와 Servlet으로 처리. 카테고리 파라미터 분기 설계",
    "2. **DB 페이징 & View 최적화**  메모리 과부하를 막는 Oracle OFFSET/FETCH 페이징 적용, 반복 JOIN을 View로 추상화",
    "3. **데이터 무결성 방어**  추천·신고 테이블에 (회원번호+게시글번호) 복합 Unique 설정으로 어뷰징 원천 차단",
    "4. **ON DELETE CASCADE 일관성**  상위 글 삭제 시 하위 댓글 등 종속 데이터 일관 자동 삭제 구조 모델링",
    "5. **중앙 집중식 보안 필터**  web.xml 및 Servlet Filter로 XSS 인젝션 방어 및 권한 우회 시도 차단",
], P2_MAIN, "12 / 16")

technical_slide(1, "03  데이터 구조", "데이터베이스 관계 및 흐름 설계", [
    "**동적 파라미터 라우팅**  Servlet 하나에서 `boardType`을 읽어 SQL 쿼리를 바인딩하고 동일한 JSP 페이지의 UI 요소를 동적 렌더링",
    "**vwPost / vwComment 뷰 작성**  자주 JOIN되는 게시글-회원-댓글 정보의 조회를 단순화하여 복잡한 DTO 매핑 코드 감소 및 SQL 성능 최적화",
    "**Filter 기반 보안**  모든 요청의 캐릭터셋 인코딩 처리 및 비인가 사용자의 마이페이지 접근 차단을 하나의 필터 체인에서 중앙 통제",
], P2_MAIN, ip("image28.png"), "스터디 데이터 모델 구조", "13 / 16")

troubleshooting_slide(1, "04  트러블슈팅", "필터 및 상태 복원 트러블슈팅", [
    "**검색·필터 후 페이지 이동 시 조건 초기화**",
    "- 단순 페이지 번호만 링크에 실어 넘기면서 검색어와 카테고리가 매번 유실되는 현상 발생.",
    "- 페이징 렌더러가 현재 브라우저의 파라미터(`category`, `keyword`)를 모두 수집하여 동적으로 페이징 링크에 빌드하도록 구현해 UX 개선.",
    "**비정상 권한 우회 시도**",
    "- 비로그인 상태에서 특정 글 수정/삭제 URL을 직접 타이핑하여 도달 시 예외 터짐.",
    "- Servlet Filter에 URI 패턴 매칭 로직을 적용하여 인가되지 않은 모든 요청을 중앙에서 1차 차단하고 로그인 화면으로 리다이렉트 처리.",
], P2_MAIN, ip("image29.png"), "상세 페이지 및 일정 바인딩 화면", "14 / 16")

retrospective_slide(1, "05  회고", "프로젝트 후기 · 회고", [
    "**협업과 Git 형상 관리**",
    "프로젝트 초기에 큰 단위로 커밋을 몰아서 처리했다가 충돌이 나 Revert했을 때 정상 코드까지 날렸습니다. 이후 '1 기능 = 1 커밋' 원칙을 철저히 실천해 리스크를 쪼개는 실무적 감각을 익혔습니다.",
    "**MVC 계층 분리의 효용**",
    "Servlet/JSP 환경에서 MVC 구조를 직접 정비하면서, 데이터베이스 스키마가 바뀌어 DAO만 고치더라도 Service나 Controller는 전혀 영향받지 않고 유지되는 SRP(단일 책임)의 위력을 깨달았습니다.",
], P2_MAIN, "15 / 16")


# ══════════════════════════════════════════════════════════════════
# CLOSING SLIDE
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, DARK)

rect(sl, 10.1, 0, 0.8, 7.5, P1_MAIN)
rect(sl, 11.0, 0, 0.8, 7.5, P2_MAIN)
rect(sl, 0.8, 2.2, 0.06, 3.2, P1_MAIN)

mt = tx(sl, 1.1, 2.0, 7.8, 4.5)
tf = mt.text_frame; tf.word_wrap = True

p1 = tf.paragraphs[0]; p1.space_after = Pt(14)
rn(p1, "THANK YOU.", 14, True, P1_MAIN)
p2 = tf.add_paragraph(); p2.space_after = Pt(24)
rn(p2, "감사합니다.", 50, True, WHITE)
p3 = tf.add_paragraph(); p3.space_after = Pt(30)
rn(p3, '"기본에 집중하고 집요하게 끝을 봅니다."', 17, color=GRAY300, italic=True)
p4 = tf.add_paragraph()
rn(p4, "010-7571-0454  ·  wlaudaoq123@naver.com  ·  github.com/wauder0430", 12.5, color=GRAY500)


# ══════════════════════════════════════════════════════════════════
# SAVE PRESENTATION
# ══════════════════════════════════════════════════════════════════
out_path = r"C:\Users\a\Desktop\쇼핑몰 프로젝트 자료모음\6_자소서, 포폴 1차 완성본_260525\포트폴리오(박지명)_v1.pptx"
prs.save(out_path)
print(f"SUCCESS: Saved portfolio to {out_path}")
print(f"Slide count: {len(prs.slides)}")
