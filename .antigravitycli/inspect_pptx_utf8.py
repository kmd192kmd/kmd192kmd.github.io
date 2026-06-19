import sys
from pptx import Presentation

pptx_path = r"C:\Users\a\Desktop\쇼핑몰 프로젝트 자료모음\6_자소서, 포폴 1차 완성본_260525\포트폴리오(김민도).pptx"
prs = Presentation(pptx_path)

output_file = r"C:\Users\a\Desktop\coding\forGithub\kmd192kmd.github.io\.antigravitycli\pptx_structure.txt"

with open(output_file, 'w', encoding='utf-8') as f:
    f.write(f"Total slides: {len(prs.slides)}\n\n")
    for idx, slide in enumerate(prs.slides):
        slide_title = ""
        if slide.shapes.title:
            slide_title = slide.shapes.title.text
        f.write(f"=== Slide {idx + 1} (Title: '{slide_title}', Layout: '{slide.slide_layout.name}') ===\n")
        for s_idx, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    f.write(f"  Shape {s_idx} [text]:\n")
                    for line in text.split('\n'):
                        f.write(f"    {line}\n")
        f.write("\n")

print("Inspection output written to pptx_structure.txt successfully!")
