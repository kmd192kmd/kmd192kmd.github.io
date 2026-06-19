import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color Palette ──
BG_DARK       = RGBColor(15, 23, 42)      # Slate-900 dark bg
BG_LIGHT      = RGBColor(255, 255, 255)   # White
TEXT_DARK     = RGBColor(15, 23, 42)      # Slate-900
TEXT_LIGHT    = RGBColor(241, 245, 249)   # Slate-100
TEXT_MUTED    = RGBColor(148, 163, 184)   # Slate-400
TEXT_SUB      = RGBColor(71, 85, 105)     # Slate-600
PRIMARY       = RGBColor(30, 58, 138)     # Blue-900
SIDEBAR_BG    = RGBColor(248, 250, 252)   # Slate-50
SIDEBAR_LINE  = RGBColor(226, 232, 240)   # Slate-200
HEADER_BG     = RGBColor(241, 245, 249)   # Slate-100 (박지명 포트 헤더 배경)

# Project accent colors
C_RED     = RGBColor(220, 38, 38)     # Gold Market
C_BLUE    = RGBColor(37, 99, 235)     # Jesiyo
C_EMERALD = RGBColor(5, 150, 105)     # DeveryTime
C_ORANGE  = RGBColor(234, 88, 12)     # 교육센터

FONT = "Malgun Gothic"

# Image paths
EXT = r"C:\Users\a\Desktop\coding\forGithub\kmd192kmd.github.io\images\extracted"
GM  = r"C:\Users\a\Desktop\coding\forGithub\kmd192kmd.github.io\images\goldmarket"

def img(folder, name):
    p = os.path.join(folder, name)
    return p if os.path.exists(p) else None

# ── Common Helpers ──

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, left, top, width, height, fill_color=None, line_color=None, line_width=0):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill_color:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_color
    else:
        shp.fill.background()
    if line_color:
        shp.line.color.rgb = line_color
        if line_width:
            shp.line.width = Pt(line_width)
    else:
        shp.line.fill.background()
    return shp

def rbox(slide, left, top, width, height, fill_color=None, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill_color:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_color
    else:
        shp.fill.background()
    if line_color:
        shp.line.color.rgb = line_color
    else:
        shp.line.fill.background()
    return shp

def tb(slide, left, top, width, height):
    return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))

def para(tf, text, font_size=12, bold=False, color=None, align=PP_ALIGN.LEFT, space_after=6, space_before=0, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.text = text
    p.font.name = FONT
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color if color else TEXT_DARK
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    return p

def add_img(slide, path, left, top, width, height, caption=None, accent=None):
    if path and os.path.exists(path):
        try:
            slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width), Inches(height))
            if caption:
                ctb = tb(slide, left, top + height + 0.05, width, 0.35)
                cp = ctb.text_frame.paragraphs[0]
                cp.text = caption
                cp.alignment = PP_ALIGN.CENTER
                cp.font.name = FONT
                cp.font.size = Pt(9)
                cp.font.italic = True
                cp.font.color.rgb = TEXT_MUTED
        except Exception as e:
            print(f"Image error {path}: {e}")

# ── Park Ji-myung Style Header (top banner on each content slide) ──
def add_header_banner(slide, project_label, accent_color):
    """박지명 스타일의 상단 배너: 왼쪽에 세로 accent bar + PORTFOLIO/포트폴리오 타이틀"""
    # 배경 헤더 영역
    h = box(slide, 0, 0, 13.333, 0.85, HEADER_BG)
    # 좌측 accent 세로 강조 bar
    box(slide, 0, 0, 0.2, 0.85, accent_color)
    # PORTFOLIO 레이블
    t1 = tb(slide, 0.3, 0.05, 2.5, 0.35)
    p1 = t1.text_frame.paragraphs[0]
    p1.text = "PORTFOLIO"
    p1.font.name = FONT
    p1.font.size = Pt(11)
    p1.font.bold = True
    p1.font.color.rgb = accent_color
    # 포트폴리오 (한글)
    t2 = tb(slide, 0.3, 0.45, 2.5, 0.35)
    p2 = t2.text_frame.paragraphs[0]
    p2.text = "포트폴리오"
    p2.font.name = FONT
    p2.font.size = Pt(9)
    p2.font.color.rgb = TEXT_MUTED
    # 프로젝트명 (오른쪽)
    t3 = tb(slide, 3.0, 0.1, 10.0, 0.65)
    p3 = t3.text_frame.paragraphs[0]
    p3.text = project_label
    p3.font.name = FONT
    p3.font.size = Pt(18)
    p3.font.bold = True
    p3.font.color.rgb = TEXT_DARK
    p3.alignment = PP_ALIGN.RIGHT

# ── Sidebar Navigation (Park Ji-myung style) ──
TABS_CONTENT = ["01 | 개요 · 기획의도", "02 | 핵심 구현", "03 | 데이터 구조", "04 | 트러블슈팅", "05 | 회고"]

def add_sidebar(slide, part_label, active_idx, accent_color):
    """5-tab vertical sidebar identical to Park Ji-myung style"""
    # Sidebar background
    sb = box(slide, 0, 0.85, 2.7, 6.65, SIDEBAR_BG, SIDEBAR_LINE)
    # Part label
    pl = tb(slide, 0.15, 0.95, 2.4, 0.45)
    p0 = pl.text_frame.paragraphs[0]
    p0.text = part_label
    p0.font.name = FONT
    p0.font.size = Pt(12)
    p0.font.bold = True
    p0.font.color.rgb = accent_color
    # Tab entries
    for idx, tab_name in enumerate(TABS_CONTENT):
        top_y = 1.55 + idx * 0.95
        is_active = (idx == active_idx)
        if is_active:
            # Active highlight background
            box(slide, 0, top_y - 0.05, 2.7, 0.88, RGBColor(255, 255, 255), SIDEBAR_LINE)
            # Active accent left bar
            box(slide, 0, top_y - 0.05, 0.18, 0.88, accent_color)
        # Tab text
        tt = tb(slide, 0.28, top_y, 2.35, 0.8)
        tt.text_frame.word_wrap = True
        tp = tt.text_frame.paragraphs[0]
        # tab number in accent color (first 2 chars)
        num, name = tab_name.split(" | ", 1)
        # Build two runs: number + name
        run1 = tp.add_run()
        run1.text = num + " | "
        run1.font.name = FONT
        run1.font.size = Pt(10)
        run1.font.bold = is_active
        run1.font.color.rgb = accent_color if is_active else TEXT_MUTED
        run2 = tp.add_run()
        run2.text = name
        run2.font.name = FONT
        run2.font.size = Pt(10)
        run2.font.bold = is_active
        run2.font.color.rgb = TEXT_DARK if is_active else TEXT_MUTED

# ── Content Area Helper ──
def content_area(slide, title, items, left=2.85, top=1.0, width=10.3, has_img=False, accent_color=None):
    """Main content area: title + bulleted items"""
    # Section title
    title_top = top
    tw = box(slide, left, title_top, width if not has_img else width * 0.52, 0.55, None)
    # accent left-border line for title
    box(slide, left, title_top + 0.05, 0.08, 0.42, accent_color if accent_color else PRIMARY)
    ttb = tb(slide, left + 0.15, title_top, (width if not has_img else width * 0.52) - 0.15, 0.55)
    ttf = ttb.text_frame
    tp_title = ttf.paragraphs[0]
    tp_title.text = title
    tp_title.font.name = FONT
    tp_title.font.size = Pt(16)
    tp_title.font.bold = True
    tp_title.font.color.rgb = TEXT_DARK

    # Items textbox
    content_w = (width * 0.50) if has_img else width
    itb = tb(slide, left, title_top + 0.62, content_w, 5.7)
    itf = itb.text_frame
    itf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = itf.paragraphs[0]
            first = False
        else:
            p = itf.add_paragraph()
        p.space_after = Pt(5)
        p.space_before = Pt(2)
        # Handle bold markers **...**
        if "**" in item:
            parts = item.split("**")
            p.text = ""
            for pi, part in enumerate(parts):
                run = p.add_run()
                run.text = part
                run.font.name = FONT
                run.font.size = Pt(11.5)
                if pi % 2 == 1:
                    run.font.bold = True
                    run.font.color.rgb = accent_color if accent_color else PRIMARY
                else:
                    run.font.color.rgb = TEXT_DARK
        else:
            p.text = item
            p.font.name = FONT
            p.font.size = Pt(11.5)
            p.font.color.rgb = TEXT_DARK

# ── Full Slide Builder ──
def content_slide(prs, project_label, part_label, active_tab, slide_title, items, accent, img_path=None, img_caption=None, img_top=None):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_bg(slide, BG_LIGHT)
    add_header_banner(slide, project_label, accent)
    add_sidebar(slide, part_label, active_tab, accent)
    # divider line between sidebar and content
    box(slide, 2.7, 0.85, 0.025, 6.65, SIDEBAR_LINE)
    has_img = img_path and os.path.exists(img_path)
    content_area(slide, slide_title, items, left=2.9, top=1.0, width=10.25, has_img=has_img, accent_color=accent)
    if has_img:
        iw = 4.9
        il = 2.9 + (10.25 * 0.50) + 0.15
        it = img_top if img_top else 1.2
        ih = min(5.8, 7.5 - it - 0.1)
        add_img(slide, img_path, il, it, iw, ih, img_caption, accent)
    return slide

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: COVER
# ══════════════════════════════════════════════════════════════════════════════
def add_cover(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_bg(slide, BG_DARK)
    # Top accent bar
    box(slide, 0, 0, 13.333, 0.12, C_RED)
    # Left accent block
    box(slide, 0, 0.12, 0.5, 7.38, RGBColor(30, 41, 59))
    # Main content
    txb = tb(slide, 1.2, 2.0, 11.0, 4.5)
    tf = txb.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = "FIRE FIGHTER TO BACKEND DEVELOPER"
    p0.font.name = FONT; p0.font.size = Pt(15); p0.font.bold = True
    p0.font.color.rgb = C_RED; p0.space_after = Pt(16)
    p1 = tf.add_paragraph()
    p1.text = "김민도 백엔드 개발자 포트폴리오"
    p1.font.name = FONT; p1.font.size = Pt(42); p1.font.bold = True
    p1.font.color.rgb = TEXT_LIGHT; p1.space_after = Pt(18)
    p2 = tf.add_paragraph()
    p2.text = '"끊임없는 질문과 집요함이 프로그램의 설계를 탄탄하게 만듭니다."'
    p2.font.name = FONT; p2.font.size = Pt(18); p2.font.color.rgb = TEXT_MUTED; p2.space_after = Pt(36)
    p3 = tf.add_paragraph()
    p3.text = "CONTACT  ·  010-3105-7821  |  kmd192kmd@naver.com  |  github.com/kmd192kmd"
    p3.font.name = FONT; p3.font.size = Pt(13); p3.font.color.rgb = TEXT_MUTED

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: PROFILE
# ══════════════════════════════════════════════════════════════════════════════
def add_profile(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_bg(slide, BG_LIGHT)
    # Header
    box(slide, 0, 0, 13.333, 0.85, HEADER_BG)
    box(slide, 0, 0, 0.2, 0.85, C_RED)
    t = tb(slide, 0.3, 0.15, 12.0, 0.55)
    p = t.text_frame.paragraphs[0]
    p.text = "ABOUT ME  &  TECHNICAL SKILLS"
    p.font.name = FONT; p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = TEXT_DARK

    # Left: Profile card
    lc = rbox(slide, 0.4, 1.0, 5.9, 6.1, SIDEBAR_BG, SIDEBAR_LINE)
    ltf = lc.text_frame; ltf.word_wrap = True
    ltf.margin_left = Inches(0.22); ltf.margin_top = Inches(0.22); ltf.margin_right = Inches(0.22)
    p_h = ltf.paragraphs[0]
    p_h.text = "Profile"
    p_h.font.name = FONT; p_h.font.size = Pt(18); p_h.font.bold = True; p_h.font.color.rgb = PRIMARY; p_h.space_after = Pt(12)
    details = [
        ("이름 / 연령", "김민도  /  1996.06.18"),
        ("학력", "전주대학교 컴퓨터공학과 졸업  (4.12 / 4.5)"),
        ("훈련 이수", "쌍용교육센터 Java Full-Stack 개발자 양성과정 (130일, 1040시간)"),
        ("자격증", "정보처리기사 (2025.12)  |  1종보통운전면허"),
        ("소방관 이력", "전직 소방관 출신. 현장의 집요함과 책임감을\n소스코드 설계와 트러블슈팅의 핵심 동력으로 삼고 있습니다."),
    ]
    for lbl, val in details:
        pl = ltf.add_paragraph()
        pl.text = f"▪ {lbl}"; pl.font.name = FONT; pl.font.size = Pt(12); pl.font.bold = True
        pl.font.color.rgb = TEXT_DARK; pl.space_before = Pt(6)
        pv = ltf.add_paragraph()
        pv.text = val; pv.font.name = FONT; pv.font.size = Pt(11)
        pv.font.color.rgb = TEXT_SUB; pv.space_after = Pt(4)

    # Right: Skills card
    rc = rbox(slide, 6.7, 1.0, 6.2, 6.1, SIDEBAR_BG, SIDEBAR_LINE)
    rtf = rc.text_frame; rtf.word_wrap = True
    rtf.margin_left = Inches(0.22); rtf.margin_top = Inches(0.22); rtf.margin_right = Inches(0.22)
    s_h = rtf.paragraphs[0]
    s_h.text = "Technical Skills"
    s_h.font.name = FONT; s_h.font.size = Pt(18); s_h.font.bold = True; s_h.font.color.rgb = PRIMARY; s_h.space_after = Pt(12)
    skills = [
        ("Back-end", "Java 21, Spring Boot 3.5, Spring Security, Spring Data JPA, QueryDSL, MyBatis"),
        ("Database & Cache", "Oracle DB, Redis (Cache/Distributed Session/Rate Limiter), Redisson 분산 락, Elasticsearch"),
        ("Infra & DevOps", "Apache Kafka, Docker, Docker Compose, GitHub Actions CI/CD, AWS EC2, Linux"),
        ("Front-end & Tools", "React, Thymeleaf, JavaScript, HTML5/CSS3, Git/GitHub, Figma, ERDCloud"),
    ]
    for cat, itms in skills:
        pc = rtf.add_paragraph()
        pc.text = f"▪ {cat}"; pc.font.name = FONT; pc.font.size = Pt(12); pc.font.bold = True
        pc.font.color.rgb = TEXT_DARK; pc.space_before = Pt(8)
        pi = rtf.add_paragraph()
        pi.text = itms; pi.font.name = FONT; pi.font.size = Pt(11)
        pi.font.color.rgb = TEXT_SUB; pi.space_after = Pt(4)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3: TABLE OF CONTENTS
# ══════════════════════════════════════════════════════════════════════════════
def add_toc(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_bg(slide, BG_LIGHT)
    box(slide, 0, 0, 13.333, 0.85, HEADER_BG)
    box(slide, 0, 0, 0.2, 0.85, TEXT_DARK)
    t = tb(slide, 0.4, 0.15, 12.0, 0.55)
    p = t.text_frame.paragraphs[0]
    p.text = "PORTFOLIO INDEX  ·  목차"
    p.font.name = FONT; p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = TEXT_DARK

    projects = [
        (C_RED,     "PART #1", "Gold Market (골드 마켓)",          "Spring Boot 기반 E-Commerce — Elasticsearch / Redis / Kafka / AWS 배포 / CI·CD"),
        (C_BLUE,    "PART #2", "Jesiyo (제시요)",                   "WebSocket·STOMP 실시간 경매 & Redis Pub/Sub — K6 부하 테스트 & Redisson 분산 락"),
        (C_EMERALD, "PART #3", "DeveryTime (데버리타임)",           "Servlet/JSP MVC Model 2 — AJAX 비동기 & DB 무결성 스터디 플랫폼"),
        (C_ORANGE,  "PART #4", "교육센터 운영 관리 시스템",         "Oracle PL/SQL · ERD 정밀 설계 — 5단계 의존성 트리 & Procedure/Trigger DB 프로젝트"),
    ]
    for i, (acc, part, name, desc) in enumerate(projects):
        top = 1.05 + i * 1.52
        # accent colored part badge
        pb = rbox(slide, 0.5, top, 1.5, 1.2, acc)
        ptf = pb.text_frame; ptf.margin_top = Inches(0.3); ptf.margin_left = Inches(0.1)
        pp = ptf.paragraphs[0]; pp.text = part
        pp.alignment = PP_ALIGN.CENTER; pp.font.name = FONT; pp.font.size = Pt(13)
        pp.font.bold = True; pp.font.color.rgb = TEXT_LIGHT
        # Content area
        dc = tb(slide, 2.3, top + 0.12, 10.6, 1.0)
        dtf = dc.text_frame; dtf.word_wrap = True
        dp0 = dtf.paragraphs[0]; dp0.text = name
        dp0.font.name = FONT; dp0.font.size = Pt(17); dp0.font.bold = True
        dp0.font.color.rgb = TEXT_DARK; dp0.space_after = Pt(4)
        dp1 = dtf.add_paragraph(); dp1.text = desc
        dp1.font.name = FONT; dp1.font.size = Pt(11.5); dp1.font.color.rgb = TEXT_SUB

# ══════════════════════════════════════════════════════════════════════════════
# DIVIDER / INTERTITLE SLIDE
# ══════════════════════════════════════════════════════════════════════════════
def add_intertitle(prs, part_num, title, subtitle, accent):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_bg(slide, BG_DARK)
    # Left vertical accent block
    box(slide, 0, 0, 0.45, 7.5, accent)
    # PART label
    pt = tb(slide, 1.0, 2.0, 11.0, 0.65)
    pp = pt.text_frame.paragraphs[0]
    pp.text = f"PART 0{part_num}."
    pp.font.name = FONT; pp.font.size = Pt(22); pp.font.bold = True; pp.font.color.rgb = accent; pp.space_after = Pt(10)
    # Main title
    tt = tb(slide, 1.0, 2.65, 11.0, 1.4)
    tp = tt.text_frame.paragraphs[0]
    tp.text = title
    tp.font.name = FONT; tp.font.size = Pt(44); tp.font.bold = True; tp.font.color.rgb = TEXT_LIGHT; tp.space_after = Pt(18)
    # Subtitle
    st = tb(slide, 1.0, 4.1, 11.0, 0.7)
    sp = st.text_frame.paragraphs[0]
    sp.text = subtitle
    sp.font.name = FONT; sp.font.size = Pt(16); sp.font.color.rgb = TEXT_MUTED
    # Tech pills row at bottom
    return slide

# ══════════════════════════════════════════════════════════════════════════════
# CLOSING SLIDE
# ══════════════════════════════════════════════════════════════════════════════
def add_closing(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_bg(slide, BG_DARK)
    box(slide, 0, 0, 13.333, 0.12, C_RED)
    t = tb(slide, 1.2, 2.2, 11.0, 4.0)
    tf = t.text_frame; tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = "FIRE FIGHTER TO BACKEND DEVELOPER"
    p0.font.name = FONT; p0.font.size = Pt(14); p0.font.bold = True
    p0.font.color.rgb = C_RED; p0.space_after = Pt(18)
    p1 = tf.add_paragraph()
    p1.text = "감사합니다!"
    p1.font.name = FONT; p1.font.size = Pt(50); p1.font.bold = True
    p1.font.color.rgb = TEXT_LIGHT; p1.space_after = Pt(20)
    p2 = tf.add_paragraph()
    p2.text = '"현장의 안전을 책임지던 소방관의 집요함으로 기초 모델링부터 동시성 제어까지 단계별로 성장해 왔습니다."'
    p2.font.name = FONT; p2.font.size = Pt(18); p2.font.color.rgb = TEXT_MUTED; p2.space_after = Pt(36)
    p3 = tf.add_paragraph()
    p3.text = "CONTACT  ·  010-3105-7821  |  kmd192kmd@naver.com  |  github.com/kmd192kmd"
    p3.font.name = FONT; p3.font.size = Pt(13); p3.font.color.rgb = TEXT_MUTED


# ══════════════════════════════════════════════════════════════════════════════
# BUILD SLIDES
# ══════════════════════════════════════════════════════════════════════════════
add_cover(prs)
add_profile(prs)
add_toc(prs)

# Short alias
def cs(project_label, part_label, tab, title, items, accent, img_name=None, img_folder=EXT, caption=None):
    ip = img(img_folder, img_name) if img_name else None
    content_slide(prs, project_label, part_label, tab, title, items, accent, ip, caption)

# ── PART 1: GOLD MARKET ──────────────────────────────────────────────────────
add_intertitle(prs, 1, "Gold Market (골드 마켓)",
               "Spring Boot · Elasticsearch · Redis · Kafka · AWS EC2 · CI/CD", C_RED)

GM_LBL = "Gold Market  |  E-Commerce 쇼핑몰 플랫폼"
GM_PART = "PART#1  골드마켓"

# 개요
cs(GM_LBL, GM_PART, 0, "개요 · 기획의도", [
    "**프로젝트 개요**: 공동구매·핫딜 서비스를 제공하는 Spring Boot 기반 E-Commerce 플랫폼",
    "**기획 의도**: 선착순 공동구매 오픈 시 대규모 트래픽 하에서의 데이터 정합성 보장, Elasticsearch/Redis 기반 조회 성능 최적화, Kafka 비동기 이벤트 격리에 집중",
    "**개발 기간**: 2026.05.19 ~ 2026.06.17 (약 4주)",
    "**개발 인원**: 4명 (백엔드 메인 프로젝트)",
    "**담당 역할**: 상품 검색/조회, 주문/반품/배송 서비스, Elasticsearch 오프로딩, GitHub Actions CI/CD 파이프라인 구축",
    "**사용 기술**: Java 21 · Spring Boot 3.5 · Spring Security · Kafka · Elasticsearch · Redis · Oracle Cloud DB · QueryDSL · React · GitHub Actions · Docker · AWS EC2",
], C_RED, "goldmarket-main.png", GM, "골드마켓 메인 화면")

# 핵심 구현 요약
cs(GM_LBL, GM_PART, 1, "핵심 구현 요약", [
    "1. **Elasticsearch 검색 오프로딩** — DB 조회를 검색 전용 엔진으로 완전 이관하여 250ms → 5~10ms 단축",
    "2. **Redis 다중 활용** — 인기 상품 캐싱(Cache Warming), 분산 세션 스토어, Rate Limiter(fixed-window) 구현",
    "3. **Kafka 비동기 이벤트 격리** — 결제 완료 이벤트를 별도 트랜잭션으로 분리하여 결제 장애 전파 차단",
    "4. **공동구매 자동화 대기열** — FIFO 기반 Spring Scheduler로 이탈자 발생 시 다음 대기자 자동 승격",
    "5. **공동구매 동시성 제어** — Redis Rate Limiter(앞단) + DB 비관적 락(최종 진입) 이중 구조",
    "6. **GitHub Actions + Docker CI/CD** — AWS EC2 RAM 8GiB 메모리 부족 문제를 GitHub Actions 빌드 분리로 해결",
    "7. **React 비동기 부분 렌더링** — 주문/취소/반품/환불 다차원 상태 UI를 공통 컴포넌트로 재사용",
], C_RED)

# 검색 최적화 상세
cs(GM_LBL, GM_PART, 1, "Elasticsearch 기반 검색 성능 최적화", [
    "**문제**: 정렬·필터 조건이 복잡해질수록 Oracle DB 직접 조회 응답시간이 250ms 이상으로 악화",
    "**1단계 — 대표 이미지 비정규화**: 매 검색 시 발생하던 이미지 JOIN 부하를 상품 테이블 컬럼 추가로 제거 (250ms)",
    "**2단계 — DB 인덱싱**: 자주 쓰이는 필터/정렬 컬럼에 Oracle 인덱스를 설정하여 응답 단축 (70ms)",
    "**3단계 — Elasticsearch 완전 이관**: 검색 조회 연산을 Elasticsearch로 오프로딩하여 DB 부하 제로화 (5~10ms 단축)",
    "**최적화 적용 기법**: Painless Script 정렬 튜닝 · 한글 형태소 Nori 분석기 · CUD 발생 시 인덱스 비동기 동기화",
], C_RED, "goldmarket-product-list.png", GM, "Elasticsearch 기반 상품 목록")

# Redis 활용
cs(GM_LBL, GM_PART, 1, "Redis 다중 활용 — 캐싱 · 세션 · Rate Limiter", [
    "**인기 상품 추천 캐싱** — 메인 추천 상품 조회 시 DB 직접 접근 차단, 서버 기동 시 캐시 워밍(Warming) + 1시간 주기 스케줄러 자동 갱신",
    "**분산 세션 스토어 구축** — 다중 서버 환경에서 로드 밸런서가 다른 서버로 연결해도 로그인 상태 유지, Spring Session Data Redis로 중앙 세션 저장소 구현",
    "**객체 직렬화 충돌 방지** — 서버 배포 시 클래스 변수 불일치로 발생하는 역직렬화 예외 차단, serialVersionUID 고정 + 직렬화 변환 방식 통일",
    "**Rate Limiter (공동구매)** — IP·계정별 단시간 반복 요청을 StringRedisTemplate fixed-window 카운터로 앞단 즉시 차단, DB까지 부하 전파 예방",
], C_RED, "goldmarket-main.png", GM, "추천 상품 캐싱 및 분산 세션")

# Kafka 비동기
cs(GM_LBL, GM_PART, 1, "Kafka 비동기 이벤트 격리 — 결제 완료 처리", [
    "**문제**: 결제 승인·주문 생성·알림 발송이 하나의 동기 트랜잭션으로 묶여 후속 처리 실패 시 결제 자체가 롤백되는 강결합 구조",
    "**1단계**: 결제 데이터를 저장하는 1차 트랜잭션을 먼저 성공 커밋 (결제 완료 확정)",
    "**2단계**: @TransactionalEventListener로 트랜잭션 커밋 완료 시점을 감지 → Apache Kafka로 이벤트 발행",
    "**3단계**: Kafka Consumer가 메시지를 수신하여 주문 생성·알림 발송을 별도 독립 트랜잭션으로 처리",
    "**효과**: 후속 기능 장애가 결제 트랜잭션으로 역류하지 않는 장애 격리 달성 — 결제 성공은 항상 보장",
], C_RED, "goldmarket-kafka.jpg", GM, "Kafka 결제 비동기 이벤트 분리도")

# DB 설계
cs(GM_LBL, GM_PART, 2, "데이터 구조 — 배송 테이블 통합 & API 최소화", [
    "**주문/발주 배타적 참조 제약조건 (Exclusive OR)**",
    "- 고객 주문(order_seq)과 업체 발주(purchase_seq) 중 한쪽만 반드시 연결되는 배타적 제약 조건 설계",
    "- 단일 배송 테이블로 양 도메인을 통합 관리 → 배송 상태 추적 쿼리 단순화",
    "**외부 API 호출 최소화 설계**",
    "- 위도/경도 최초 1회 저장: 주소 등록 시 카카오맵 API로 한 번만 조회하여 DB에 영구 저장, 이후 내부 연산만 수행",
    "- 공휴일 테이블 구축: 배송 예정일 계산 시 매번 외부 API 호출 대신 한 달 1회 스케줄러로 동기화하여 DB에서 처리",
], C_RED, "goldmarket-delivery-erd.png", GM, "배송 ERD 및 주문·발주 통합 모델")

# React + 뒤로가기
cs(GM_LBL, GM_PART, 3, "트러블슈팅 — React 부분 렌더링 & 뒤로가기 캐시 문제", [
    "**React 비동기 부분 렌더링 도입**",
    "- 주문 목록·배송 조회·취소/반품/교환/환불 내역에서 동일한 UI 패턴이 반복되어 React 컴포넌트화로 코드 재사용성 극대화",
    "- 주문 상태·환불 승인 여부 등 다양한 조건에 따른 분기 렌더링을 State/Props로 처리, JSON API 비동기 통신 연동",
    "**뒤로가기 시 최근 본 상품 갱신 미반영 문제** (Thymeleaf 일반 페이지 영역)",
    "- 브라우저 뒤로가기 기본 동작(캐시 복원)을 History API로 가로채어 정지",
    "- pageshow 이벤트를 감지하여 최근 본 상품 영역만 Fetch API로 비동기 재호출 → 즉각적인 데이터 일관성 달성",
], C_RED, "goldmarket-orders-delivery.png", GM, "React 기반 주문/배송 추적 UI")

# 공동구매
cs(GM_LBL, GM_PART, 3, "트러블슈팅 — 공동구매 자동 대기열 & 동시성 제어", [
    "**공동구매 자동 참여 전환 대기열 구축**",
    "- 선착순 참여자 이탈 시 FIFO 대기열 첫 번째 대기자에게 24시간 결제 기회 자동 부여",
    "- 기한 초과 시 자동 회수 후 다음 대기자에게 이전 — Spring Scheduler + JPA 트랜잭션 전파로 구현",
    "**공동구매 동시 참여 이중 보안 구조**",
    "- Rate Limiter (앞단): Redis StringRedisTemplate fixed-window 카운터로 IP·계정별 반복 요청 즉시 차단",
    "- 비관적 락 (최종 진입): DB Pessimistic Lock으로 동시 진입 트랜잭션을 순차 처리하여 재고 초과 판매 원천 방지",
], C_RED, "goldmarket-groupbuy.png", GM, "공동구매 FIFO 대기열 및 동시성 제어")

# 주문 취소 + Kafka
cs(GM_LBL, GM_PART, 3, "트러블슈팅 — 주문 취소 동시성 & 환불 멱등성 보장", [
    "**문제**: 연속 클릭으로 중복 환불 위험 + Spring Batch 배송 배치 작업과 취소 요청 충돌",
    "**화면 단**: 취소 버튼 즉시 disabled 처리로 물리적 중복 요청 차단",
    "**DB 비관적 락 (Pessimistic Lock)**: 백엔드 진입 즉시 락 획득 → 충돌 요청 순차 처리",
    "**영속성 컨텍스트 최신화**: Spring Batch가 반영한 최신 배송 상태를 읽기 위해 EntityManager.refresh()로 강제 동기화",
    "**외부 결제 API 멱등성 보장**: 이미 취소된 결제 건(ALREADY_CANCELED_PAYMENT)은 성공으로 흡수하여 정상 리턴",
], C_RED, "goldmarket-cancel.png", GM, "주문 취소 완료 화면")

# CI/CD
cs(GM_LBL, GM_PART, 3, "트러블슈팅 — GitHub Actions CI/CD & AWS EC2 배포", [
    "**문제**: RAM 8GiB 제한 AWS EC2 서버에서 직접 Gradle 빌드 시 메모리 부족으로 전체 서비스 중단",
    "**해결**: 빌드를 GitHub Actions Runner 환경으로 외부 분리",
    "- GitHub Actions에서 JAR 빌드 → Docker 이미지 생성 → Docker Hub에 푸시",
    "- EC2에서는 이미지를 Pull하여 Docker Compose로 컨테이너 기동 → 서버 메모리 빌드 부하 제로화",
    "**보안 사고 경험**: Git에 외부 연동용 API Key 노출 → git-filter-repo로 Git History 수정 + 환경 변수 .env 격리로 재발 방지",
    "**자동화 배포 파이프라인**: main 브랜치 Push 시 Actions 자동 실행 → Docker Hub 업로드 → EC2 SSH 원격 접속 → Compose 재시작",
], C_RED, "goldmarket-cicd.jpg", GM, "GitHub Actions & AWS EC2 CI/CD 구성")

# 매출 최적화
cs(GM_LBL, GM_PART, 3, "트러블슈팅 — 매출 통계 N+1 최적화 & 대량 쿠폰 발급", [
    "**N+1 쿼리 병목**: 매출 대시보드 조회 시 결제·상품 정보를 개별 요청으로 가져오는 다중 쿼리 발생",
    "**해결**: QueryDSL Dynamic 조건문을 단일 JOIN 쿼리로 통합 처리 → 조회 속도 개선",
    "**정밀 소수 연산**: 금융 정산의 소수점 오차 예방을 위해 double 대신 BigDecimal 타입 사용",
    "**대량 쿠폰 발급 병목**: 수천 건 쿠폰 단건 INSERT 루프로 극심한 지연 발생",
    "**해결**: JDBC Batch Insert 방식 적용 → 한 번의 배치 실행으로 대량 등록 속도 획기적 향상",
], C_RED, "goldmarket-admin-stats.png", GM, "관리자 매출 통계 대시보드")

# 회고
cs(GM_LBL, GM_PART, 4, "프로젝트 완수 후기 · 회고", [
    "**배포 및 보안 경험**",
    "GitHub Actions CI/CD 구축 중 메모리 부족 문제와 AWS EC2 배포 오류들을 하나씩 해결하며 서버·배포에 대한 실전 경험을 쌓았습니다. 특히 Git에 API Key가 노출되는 사고를 겪으며 Git History 수정 및 환경 변수 격리의 중요성을 직접 체험했습니다.",
    "**바이브 코딩 시대와 트레이드 오프 설계 능력**",
    "AI가 코드를 빠르게 완성해 주는 시대가 되었지만, 시스템 구조 설계·검색 성능 최적화 시 어느 수준까지 고도화할지, 비용 대비 효과를 어떻게 결정할지에 대한 트레이드 오프 판단은 결국 개발자의 몫입니다.",
    "인프라와 분산 시스템 등 다양한 기술적 깊이를 이해하고 트레이드 오프를 고민할 줄 아는 개발자로 성장하겠습니다.",
], C_RED, "후기.png", GM, "프로젝트 완수 후기")


# ── PART 2: JESIYO ───────────────────────────────────────────────────────────
add_intertitle(prs, 2, "Jesiyo (제시요)",
               "WebSocket · STOMP · Redis Pub/Sub — 실시간 경매 & Redisson 분산 락 · K6 부하 테스트", C_BLUE)

JS_LBL  = "Jesiyo  |  실시간 경매·중고거래 플랫폼"
JS_PART = "PART#2  제시요"

cs(JS_LBL, JS_PART, 0, "개요 · 기획의도", [
    "**프로젝트 개요**: 실시간 비동기 통신과 동시성 제어를 고려한 경매·중고거래 플랫폼",
    "**기획 의도**: 채팅과 경매 입찰이 분리되어 있던 UX를 채팅방 내 경매 탭으로 통합. 대용량 입찰 트래픽에서 최고가 정합성 보장 및 K6 부하 테스트로 시스템 한계 정밀 검증",
    "**개발 기간**: 2026.04.15 ~ 2026.04.29 (15일)",
    "**개발 인원**: 4명",
    "**담당 역할**: 경매·실시간 입찰 도메인 전담, Redisson 분산 락 동시성 제어, Redis 인메모리 캐싱, K6 부하 테스트 수치 도출",
    "**사용 기술**: Oracle 19c · Tomcat 9 · Java 11 · Spring 5.3 · Spring Security · MyBatis · Redis · Redisson · JUnit4 · K6",
], C_BLUE, "image31.png", EXT, "제시요 실시간 경매 화면")

cs(JS_LBL, JS_PART, 1, "핵심 구현 사항", [
    "1. **WebSocket/STOMP 실시간 입찰·채팅 통합**: 채팅방 안에 경매 탭 배치 → 사용자가 채팅하며 동시에 경매 참여, 화면 전환 오버헤드 제거",
    "2. **Redisson 분산 락 동시성 제어**: 대량 동시 입찰 시 최고가 로직 무결성 보장 — Redis 기반 분산 락으로 스레드 순차 제어",
    "3. **Redis 최고가 캐시 (Cache-Aside)**: 락 대기 병목을 DB가 아닌 Redis 인메모리 캐시 조회로 상쇄하여 응답 속도 향상",
    "4. **클라이언트 단 타이머 렌더링**: 남은 입찰 시간 표시를 서버 Polling 대신 브라우저 setInterval로 독립 계산 → 서버 트래픽 부하 절감",
    "5. **(auction_seq, status) 복합 인덱스**: 최고가 갱신 시 이전 최고가 상태 변경 쿼리에 복합 인덱스 적용 → Oracle DB 스캔 성능 확보",
], C_BLUE)

cs(JS_LBL, JS_PART, 2, "데이터 구조 — 시스템 아키텍처 & DB 설계", [
    "**시스템 아키텍처**: HTTP 요청과 WebSocket/STOMP 프로토콜을 분리 처리. 다중 서버 확장 시 실시간 메시지 중계를 위해 Redis Pub/Sub 구조 적용",
    "**슈퍼타입 테이블 모델**: auction_master 전역 PK를 상속하는 슈퍼타입-서브타입 구조로 일반 경매·라이브 경매 공통 연관 테이블 단순화. OCP 원칙 구현 (신규 경매 형태 추가 시 기존 구조 영향 최소화)",
    "**입찰기록 테이블 & 복합 인덱스**: 회원-경매 N:N 관계 해소, @Transactional 양쪽 쿼리 원자적 처리, (auction_seq, status) 복합 인덱스로 상태 변경 조회 최적화",
], C_BLUE, "image33.png", EXT, "시스템 아키텍처 설계도")

cs(JS_LBL, JS_PART, 3, "트러블슈팅 — K6 부하 테스트 & 성능 개선", [
    "**K6 테스트 1 — 동시성 중복 입찰 이슈 (Redisson 분산 락 도입)**",
    "- 1,000명 동시 입찰 부하에서 최고가 로직이 무시되어 동일 가격 중복 입찰 발생",
    "- Redisson 분산 락 도입으로 동시 접근 스레드 순차 제어 → K-5000 환경에서도 데이터 무결성 확보",
    "**K6 테스트 2 — 분산 락 병목 극복 (Redis 인메모리 캐시 전환)**",
    "- 락 도입 후 대기열 증가로 응답 성공률 저하 (92.93%)",
    "- 최고가 검증 로직을 DB 디스크 → Redis 인메모리 캐시 조회로 전환",
    "- 결과: 응답 성공률 98.61% 회복 · 평균 응답속도 **794ms → 2ms (390배 단축)** · TPS **1,061 → 9,593 (9배 증가)**",
], C_BLUE, "image37.png", EXT, "K6 부하 테스트 성능 지표")

cs(JS_LBL, JS_PART, 4, "프로젝트 후기 · 회고", [
    "**동시성 트레이드 오프에 대한 통찰**",
    "단순 기능 구현이 아닌 K6 실측 테스트 기반으로 병목을 분석하고 Redisson 락 + Redis 캐시의 최적 결합 구조를 도출한 실전적 최적화 경험을 얻었습니다.",
    "**메시지 큐의 필요성 체감**",
    "Redis Pub/Sub으로 실시간 브로드캐스팅 확장성을 경험했으나, 완벽한 장애 격리와 비동기 보장을 위해 Apache Kafka 같은 전문 분산 메시징 브로커가 필요함을 깨달았고 — 이 경험이 Gold Market에서 Kafka 직접 도입으로 이어졌습니다.",
    "**AI 도구의 비판적 수용**",
    "Redis 뼈대 구성·프론트 타이머는 AI 도구로 신속히 완성했으나, 동시성 오류 지점 규명과 슈퍼타입 설계의 장단점 판단은 개발자의 고유 역량임을 실감했습니다.",
], C_BLUE)


# ── PART 3: DEVERY TIME ──────────────────────────────────────────────────────
add_intertitle(prs, 3, "DeveryTime (데버리타임)",
               "Servlet/JSP MVC Model 2 — AJAX 비동기 & DB 무결성 학습관리 플랫폼", C_EMERALD)

DT_LBL  = "DeveryTime  |  학습 관리 사이트"
DT_PART = "PART#3  데버리타임"

cs(DT_LBL, DT_PART, 0, "개요 · 기획의도", [
    "**프로젝트 개요**: Servlet/JSP 환경에서 MVC Model 2 구조로 설계한 스터디·학습관리 게시판 웹 사이트",
    "**기획 의도**: 프레임워크(Spring) 없이 순수 서블릿 구조와 JSP Model 2를 밑바닥부터 구현하여 HTTP 데이터 흐름과 MVC 패턴을 깊이 이해하고자 기획",
    "**개발 기간**: 2026.03.19 ~ 2026.03.31 (13일)",
    "**개발 인원**: 4명",
    "**담당 역할**: 스터디 도메인·다중 게시판 관리, 권한/인증 Filter 구현, AJAX 비동기 To-Do, FullCalendar.js 커스터마이징",
    "**사용 기술**: Oracle 21c · Tomcat 9 · Java 21 · Servlet/JSP (Model 2) · JDBC · FullCalendar.js · HTML5/CSS3/JavaScript",
], C_EMERALD, "image20.png", EXT, "데버리타임 메인 화면")

cs(DT_LBL, DT_PART, 1, "핵심 구현 사항", [
    "1. **JSP Model 2 MVC 아키텍처**: Controller(요청 흐름) / Service(비즈니스 로직) / DAO(DB 접근) 계층 분리 — 단일 책임 원칙(SRP) 준수",
    "2. **JSTL 권한별 UI 분기**: 스터디 상태(모집중/참여중)·가입 권한(방장/일반)에 따른 버튼 노출 제어 + 서블릿 Filter 인가 방어",
    "3. **FullCalendar.js 데이터 바인딩**: Java DTO → JS 이벤트 객체 배열 변환 매핑, CSS 오버라이딩으로 다중 일정 UI 가독성 개선",
    "4. **Fetch API 비동기 To-Do & 롤백**: 서버 응답 성공 후에만 UI 완료 처리, 실패 시 체크박스 즉시 복원 — 클라이언트·서버 상태 정합성 보장",
    "5. **데이터 무결성 설계**: N:N 관계 해소 매핑 테이블 + (member_seq, study_seq) Unique 제약 + ON DELETE CASCADE + Soft Delete 대응 트리거",
], C_EMERALD)

cs(DT_LBL, DT_PART, 2, "데이터 구조 — 무결성 & 관계 모델링", [
    "**N:N 관계 해소**: 회원(Member)·스터디(Study) 다대다 관계를 study_member 연결 테이블로 분리",
    "**중복 가입 방지**: (member_seq, study_seq) 복합 Unique 제약 조건 → DB 구조적 레벨에서 이중 가입 완전 차단",
    "**ON DELETE CASCADE**: 스터디 삭제 시 관련 To-Do·멤버 데이터 고아 방지",
    "**Soft Delete 대응 트리거**: 회원 탈퇴(status 변경)가 물리적 삭제가 아니므로 CASCADE 미발동 → PL/SQL 트리거로 탈퇴 감지 시 study_member 자동 DELETE",
    "**트리거 vs 애플리케이션 로직 트레이드오프**: 트리거는 편리하지만 비즈니스 로직이 DB에 은닉되는 문제 인식 → 리팩토링 시 Java Service 레이어로 이관 가이드라인 수립",
], C_EMERALD, "image28.png", EXT, "스터디 테이블 구조 설계")

cs(DT_LBL, DT_PART, 3, "트러블슈팅 — 비정규화 성능 최적화 & 논리적 삭제 대응", [
    "**정원 초과 방지 복합 트리거 → 비정규화 대안 설계**",
    "- 기존: 가입 시마다 전체 인원을 COUNT하는 트리거 → 대량 더미데이터 삽입 시 성능 저하 우려",
    "- 개선: 스터디 테이블에 '현재 인원수' 컬럼 추가 (비정규화) + CHECK (현재인원 <= 모집인원) 제약조건",
    "- 효과: 가입 트랜잭션마다 단순 ±1 UPDATE만 수행하여 오버헤드 최소화",
    "**Soft Delete 회원 탈퇴 처리**",
    "- 탈퇴 후 study_member에 유령 데이터 잔존 문제를 PL/SQL 트리거로 임시 해결",
    "- 장기 개선 방향: Java 애플리케이션 Service 레이어 내 트랜잭션 순차 처리로 이관하여 테스트 가능성·투명성 확보",
], C_EMERALD, "image29.png", EXT, "회원탈퇴 종속 데이터 삭제 트리거")

cs(DT_LBL, DT_PART, 4, "프로젝트 후기 · 회고", [
    "**순수 자바 웹 아키텍처 이해**",
    "Spring 없이 Servlet 라이프사이클을 직접 구현하며 HTTP Request/Response 흐름과 서블릿 Filter 보안 동작 방식을 하부 구조부터 체득했습니다.",
    "**협업 규칙과 커밋 관리의 뼈아픈 교훈**",
    "프로젝트 초기 명명 규칙과 파라미터 규격을 확실히 정하지 않아 병합 과정에서 극심한 디버깅 오버헤드를 겪었습니다. 이후 '1 기능 = 1 커밋' 원칙과 기능별 브랜치 전략을 적극 실천하게 된 계기가 되었습니다.",
], C_EMERALD)


# ── PART 4: 교육센터 ──────────────────────────────────────────────────────────
add_intertitle(prs, 4, "교육센터 운영 관리 시스템",
               "Oracle PL/SQL · ERD 정밀 설계 — 5단계 의존성 트리 & Procedure/Trigger DB 프로젝트", C_ORANGE)

ST_LBL  = "교육센터 운영관리  |  DB 모델링 프로젝트"
ST_PART = "PART#4  교육센터"

cs(ST_LBL, ST_PART, 0, "개요 · 기획의도", [
    "**프로젝트 개요**: 웹 구현 이전, RDBMS 핵심 기본기를 습득하기 위한 집중 데이터 설계·DB 프로젝트",
    "**기획 의도**: 요구분석서·기능 명세서·ERD 설계서·DDL/DML 정의서·PL/SQL Procedure/Trigger 작성까지 RDBMS 산출물 전 과정 직접 이행",
    "**개발 기간**: 2026.02.03 ~ 2026.02.10 (8일)",
    "**개발 인원**: 4명",
    "**담당 역할**: 전체 테이블 설계·무결성 제약 정의, 5단계 의존성 트리 수립, 정원 초과 방지 트리거·프로시저 개발",
    "**사용 기술**: Oracle 21c · Oracle SQL Developer · PL/SQL · ERDCloud · Figma",
], C_ORANGE, "image6.png", EXT, "교육센터 요구분석서 산출물")

cs(ST_LBL, ST_PART, 1, "핵심 구현 사항", [
    "1. **5단계 테이블 깊이 설계도**: FK 의존성에 따라 부모→자식 순서로 5단계 분류, 더미 데이터 순차 삽입 시 FK 위배 에러 원천 방지",
    "2. **관계 모델링**: 1:N, N:N 해소, 식별/비식별 관계 명확 분류 후 ERD 완성",
    "3. **정규화 vs 반정규화 설계**: 읽기 트래픽 패턴 분석 후 JOIN 최소화를 위한 반정규화 부분 적용 (필기/실기 개별 컬럼 분리)",
    "4. **CHECK 제약조건 활용**: 배점 합 100점 초과 방지를 트리거 대신 컬럼 레벨 CHECK 제약으로 처리 → 오버헤드 최소화",
    "5. **수강생 등록 정원 초과 방지 트리거**: BEFORE ROW 트리거로 정원 초과 시 RAISE_APPLICATION_ERROR 즉시 발생",
], C_ORANGE)

cs(ST_LBL, ST_PART, 2, "데이터 구조 — 테이블 의존성 트리 & 관계 재정의", [
    "**5단계 의존성 트리 설계도**: 30여 개 테이블을 FK 참조 방향에 따라 5단계로 분류하여 데이터 입력 순서 충돌 제거",
    "**취업성공수당 관계 재정의**",
    "- 초기 설계: 수당 테이블 → 취업상태 테이블 종속 (일회성 지급 제약 구현 어려움)",
    "- 개선: 수강생 테이블과 1:1 식별 관계로 직접 연결 → DB 구조적 레벨에서 일회성 지급 제약 달성",
    "**비정규화 설계 (시험 배점 테이블)**: 필기/실기 자주 동시 조회됨 → 세로형(정규화) 대신 가로형 컬럼 분리(반정규화) 채택",
    "**nullable FK 설계**: 강의실 미배정 교구 컴퓨터 존재 → 컴퓨터 테이블 강의실 FK NOT NULL 제약 해제",
], C_ORANGE, "image8.png", EXT, "전체 테이블 의존성 트리 설계도")

cs(ST_LBL, ST_PART, 3, "트러블슈팅 — 제약 조건 최적화 & 트리거 설계", [
    "**CHECK 제약조건으로 트리거 대체**",
    "- 배점 합 검증 로직에 트리거 사용 시 컴파일 부하·숨겨진 로직 문제 발생",
    "- 컬럼 레벨 CHECK (필기 + 실기 + 출석 <= 100) 제약조건으로 변경 → 가볍고 직관적인 검증 달성",
    "**수강생 등록 정원 초과 방지 트리거 설계**",
    "- BEFORE ROW 트리거: INSERT 직전 강의실 정원·현재 등록 인원 실시간 비교",
    "- 정원 초과 시 RAISE_APPLICATION_ERROR로 트랜잭션 중단 → 유해 데이터 삽입 완전 차단",
    "**트리거 남용 제거 원칙**: 단순 검증은 CHECK 제약으로, 복잡한 비즈니스 로직은 Java Service 레이어로 처리하는 가이드라인 수립",
], C_ORANGE, "image10.png", EXT, "정원 초과 방지 트리거 코드")

cs(ST_LBL, ST_PART, 4, "프로젝트 후기 · 회고", [
    "**데이터 설계의 주춧돌**",
    "ERD 모델링과 정규화/반정규화 트레이드오프를 실질적으로 검증하며 단순 SQL 작성자가 아닌 정교한 DB 스키마 설계 역량의 필요성을 실감했습니다.",
    "**RDBMS 제약 조건의 신뢰성**",
    "PK/FK, CHECK, UNIQUE 제약 조건으로 데이터 최종 방어선을 DB 레벨에서 구축하는 법을 배워, 애플리케이션 단에만 의존하지 않는 데이터 일관성 확보 역량을 다졌습니다.",
    "**점진적·반복적 모델링 실천**",
    "작은 도메인 단위부터 테이블 배치·의존성 부여 → 테스트 병행의 반복적 모델링 기법의 중요성을 깨달았습니다.",
], C_ORANGE)


# ──────────────────────────────────────────────────────────────────────────────
add_closing(prs)

out = r"C:\Users\a\Desktop\쇼핑몰 프로젝트 자료모음\6_자소서, 포폴 1차 완성본_260525\포트폴리오(김민도)_v2.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Total slides: {len(prs.slides)}")
