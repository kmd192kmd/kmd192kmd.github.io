from pptx import Presentation

def reorder_slides(prs, new_order):
    # prs.slides._sldIdLst is the XML element containing slide IDs
    sldIdLst = prs.slides._sldIdLst
    slides_elements = list(sldIdLst)
    
    # Clear list
    for el in slides_elements:
        sldIdLst.remove(el)
        
    # Append in new order
    for idx in new_order:
        sldIdLst.append(slides_elements[idx])

# Test load and print slide count
pptx_path = r"C:\Users\a\Desktop\쇼핑몰 프로젝트 자료모음\6_자소서, 포폴 1차 완성본_260525\포트폴리오(김민도).pptx"
prs = Presentation(pptx_path)
print("Original slide count:", len(prs.slides))
