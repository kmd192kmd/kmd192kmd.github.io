from pptx import Presentation

pptx_path = r"C:\Users\a\Desktop\쇼핑몰 프로젝트 자료모음\6_자소서, 포폴 1차 완성본_260525\포트폴리오(김민도).pptx"

prs = Presentation(pptx_path)
print(f"Number of slides: {len(prs.slides)}")

for idx, slide in enumerate(prs.slides):
    slide_title = ""
    if slide.shapes.title:
        slide_title = slide.shapes.title.text
    print(f"Slide {idx + 1}: Title='{slide_title}'")
    
    # Print slide layout name
    print(f"  Layout: {slide.slide_layout.name}")
    
    # Print shapes and their text if small
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text.strip().replace('\n', ' ')
            if text:
                print(f"    Text: {text[:100]}...")
