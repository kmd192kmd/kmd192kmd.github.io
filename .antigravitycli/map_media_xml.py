import os
from pptx import Presentation
from pptx.oxml.ns import qn

pptx_path = r"C:\Users\a\Desktop\쇼핑몰 프로젝트 자료모음\6_자소서, 포폴 1차 완성본_260525\포트폴리오(김민도).pptx"
prs = Presentation(pptx_path)

output_file = r"C:\Users\a\Desktop\coding\forGithub\kmd192kmd.github.io\.antigravitycli\slide_media_xml_mapping.txt"

with open(output_file, 'w', encoding='utf-8') as out_f:
    for idx, slide in enumerate(prs.slides):
        out_f.write(f"=== Slide {idx + 1} ===\n")
        
        # We can extract the relationships of the slide to media files
        rels = slide.part.rels
        rId_to_target = {}
        for rId, rel in rels.items():
            if "media" in rel.target_ref:
                rId_to_target[rId] = os.path.basename(rel.target_ref)
                
        # Now let's traverse the slide shapes and find the rId
        for s_idx, shape in enumerate(slide.shapes):
            shape_xml = shape._element
            # Find any blip elements (which point to images)
            blips = shape_xml.findall('.//' + qn('a:blip'))
            for b_idx, blip in enumerate(blips):
                embed_rId = blip.get(qn('r:embed'))
                target_img = rId_to_target.get(embed_rId, embed_rId)
                out_f.write(f"  Shape {s_idx} ('{shape.name}', type={shape.shape_type}) -> Blip {b_idx}: {target_img}\n")

import os
print("XML media mapping complete!")
