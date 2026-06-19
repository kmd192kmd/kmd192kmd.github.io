from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

pptx_path = r"C:\Users\a\Desktop\쇼핑몰 프로젝트 자료모음\6_자소서, 포폴 1차 완성본_260525\포트폴리오(김민도).pptx"
prs = Presentation(pptx_path)

slide = prs.slides[2] # Slide 3 (0-indexed 2)
print(f"Slide 3 has {len(slide.shapes)} shapes:")

for idx, shape in enumerate(slide.shapes):
    print(f"Shape {idx}: '{shape.name}', type={shape.shape_type}")
    print(f"  Position: left={shape.left}, top={shape.top}, width={shape.width}, height={shape.height}")
    if shape.has_text_frame:
        print(f"  Text: {shape.text.strip()}")
