import os
import hashlib
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

pptx_path = r"C:\Users\a\Desktop\쇼핑몰 프로젝트 자료모음\6_자소서, 포폴 1차 완성본_260525\포트폴리오(김민도).pptx"
extract_dir = r"C:\Users\a\Desktop\coding\forGithub\kmd192kmd.github.io\images\extracted"

# Calculate MD5 for all extracted files
md5_to_filename = {}
for file in os.listdir(extract_dir):
    filepath = os.path.join(extract_dir, file)
    with open(filepath, 'rb') as f:
        data = f.read()
        md5_val = hashlib.md5(data).hexdigest()
        md5_to_filename[md5_val] = file

prs = Presentation(pptx_path)

output_file = r"C:\Users\a\Desktop\coding\forGithub\kmd192kmd.github.io\.antigravitycli\slide_image_md5_mapping.txt"

with open(output_file, 'w', encoding='utf-8') as out_f:
    for idx, slide in enumerate(prs.slides):
        out_f.write(f"=== Slide {idx + 1} ===\n")
        title = ""
        # Find any text that looks like a title
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                lines = shape.text.strip().split('\n')
                if len(lines[0]) < 50:
                    title = lines[0]
                    break
        out_f.write(f"  Title search: '{title}'\n")
        
        # Helper to check shape fill
        def check_fill(shape, path_name):
            if shape.fill and hasattr(shape.fill, 'type') and shape.fill.type == 6: # PICTURE
                try:
                    img_blob = shape.fill.image.blob
                    img_md5 = hashlib.md5(img_blob).hexdigest()
                    filename = md5_to_filename.get(img_md5, "unknown")
                    out_f.write(f"    {path_name} (Fill Picture): filename={filename}, name='{shape.name}', size={shape.fill.image.size}\n")
                except Exception as e:
                    out_f.write(f"    {path_name} (Fill Picture Error): {str(e)}\n")
        
        for s_idx, shape in enumerate(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img_blob = shape.image.blob
                    img_md5 = hashlib.md5(img_blob).hexdigest()
                    filename = md5_to_filename.get(img_md5, "unknown")
                    out_f.write(f"    Shape {s_idx} [PICTURE]: filename={filename}, name='{shape.name}', size={shape.image.size}\n")
                except Exception as e:
                    out_f.write(f"    Shape {s_idx} [PICTURE Error]: {str(e)}\n")
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for sub_idx, sub_shape in enumerate(shape.shapes):
                    if sub_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            img_blob = sub_shape.image.blob
                            img_md5 = hashlib.md5(img_blob).hexdigest()
                            filename = md5_to_filename.get(img_md5, "unknown")
                            out_f.write(f"    Shape {s_idx} Group Sub {sub_idx} [PICTURE]: filename={filename}, name='{sub_shape.name}'\n")
                        except Exception as e:
                            pass
                    check_fill(sub_shape, f"Shape {s_idx} Group Sub {sub_idx}")
            else:
                check_fill(shape, f"Shape {s_idx}")
        out_f.write("\n")

print("MD5 image mapping complete!")
