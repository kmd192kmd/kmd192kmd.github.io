"""
포트폴리오(김민도)_v10.pptx - 김민도 포트폴리오 (구성 B: 메인 프로젝트 집중 Deep-Dive 구성)
Composition: "Core Project Deep-Dive & Highlight"
- 구성 구조:
  - Part 1: 메인 프로젝트 집중 분석 (Gold Market) - 12개 슬라이드에 걸친 상세 구현 및 검색/캐싱/이벤트 격리 최적화 Deep-Dive
  - Part 2: 기초 역량 프로젝트 요약 (Jesiyo, DeveryTime, 교육센터) - 각 프로젝트당 1장의 고압축 요약 슬라이드로 신속 요약
- 기획의도: 시니어 면접관들이 선호하는 이력서 구성 방식으로, 가장 난이도가 높고 최신 스택을 활용한 메인 프로젝트(골드 마켓)에 발표 역량의 80%를 집중하고, 나머지 프로젝트는 백엔드 기본기(동시성, MVC, DB 모델링)를 보여주는 배경 정보로 다루어 전문성과 설득력을 극대화합니다.
- 디자인: 딥 슬레이트 다크 (slate-900) 배경 + 차콜 카드 (slate-800) + 네온 브랜드 컬러 악센트의 프리미엄 다크 테크 모드
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── DESIGN TOKENS (Dark Mode) ───
DARK         = RGBColor(15, 23, 42)      # Deep Slate (slate-900)
CARD_BG      = RGBColor(30, 41, 59)      # slate-800
CARD_BORDER  = RGBColor(51, 65, 85)      # slate-700
WHITE        = RGBColor(255, 255, 255)   # White
GRAY300      = RGBColor(203, 213, 225)   # slate-300
GRAY400      = RGBColor(148, 163, 184)   # slate-400
GRAY500      = RGBColor(100, 116, 139)   # slate-500

# Brand Colors (Neon Accent)
P1_MAIN  = RGBColor(244, 63, 94)     # Neon Rose/Crimson (Gold Market)
P1_LINE  = RGBColor(253, 164, 186)   # Rose-300

P2_MAIN  = RGBColor(168, 85, 247)    # Neon Violet (Jesiyo)
P2_LINE  = RGBColor(216, 180, 254)   # Violet-300

P3_MAIN  = RGBColor(16, 185, 129)    # Neon Emerald (DeveryTime)
P3_LINE  = RGBColor(167, 243, 208)   # Emerald-300

P4_MAIN  = RGBColor(245, 158, 11)    # Neon Amber (교육센터)
P4_LINE  = RGBColor(253, 201, 138)   # Amber-300

FN = "Malgun Gothic"

# Image dirs
EXT = r"C:\Users\a\Desktop\coding\forGithub\kmd192kmd.github.io\images\extracted"
GMD = r"C:\Users\a\Desktop\coding\forGithub\kmd192kmd.github.io\images\goldmarket"

def ip(folder, name):
    p = os.path.join(folder, name)
    return p if os.path.exists(p) else None

# ─── PRIMITIVES ───
def bg(slide, c):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = c

def rect(slide, l, t, w, h, fill=None, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
    else:
        s.line.fill.background()
    return s

def rrect(slide, l, t, w, h, fill=None, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
    else:
        s.line.fill.background()
    return s

def tx(slide, l, t, w, h):
    return slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))

def rn(p, text, size=12, bold=False, color=WHITE, italic=False):
    r = p.add_run(); r.text = text; r.font.name = FN
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.italic = italic
    return r

def pic(slide, path, l, t, w, h, cap=None):
    if path and os.path.exists(path):
        try:
            slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
            if cap:
                tb = tx(slide, l, t + h + 0.05, w, 0.3)
                p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
                rn(p, cap, 9, italic=True, color=GRAY400)
        except Exception as e:
            print(f"Picture error: {e}")

# ─── TECH PILLS DRAW FOR DIVIDERS ───
def draw_tech_pills(slide, start_x, start_y, pills_list, fill_color, text_color):
    curr_x = start_x
    curr_y = start_y
    for pill in pills_list:
        pill_w = len(pill) * 0.085 + 0.4
        if curr_x + pill_w > 12.5:
            curr_x = start_x
            curr_y += 0.45
        
        p_shp = rrect(slide, curr_x, curr_y, pill_w, 0.35, DARK, text_color)
        tf = p_shp.text_frame
        tf.word_wrap = False
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.03)
        tf.margin_bottom = Inches(0.03)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        rn(p, pill, 9.5, True, text_color)
        curr_x += pill_w + 0.12

# ─── MINIMAL CLEAN HEADER ───
PROJ_NAMES = ["Gold Market", "Foundational Projects"]

def add_clean_header(slide, active_proj_idx, section_name, title, accent):
    rect(slide, 0.8, 0.45, 0.08, 0.65, accent)
    tb = tx(slide, 1.05, 0.35, 8.5, 0.85)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    rn(p, f"{section_name.upper()}   //\n", 9.5, True, accent)
    rn(p, title, 20, True, WHITE)
    
    tb_r = tx(slide, 9.5, 0.45, 3.033, 0.65)
    tf_r = tb_r.text_frame
    tf_r.margin_left = tf_r.margin_top = tf_r.margin_right = tf_r.margin_bottom = 0
    p_r = tf_r.paragraphs[0]; p_r.alignment = PP_ALIGN.RIGHT
    rn(p_r, PROJ_NAMES[active_proj_idx].upper(), 11, True, accent)
    rn(p_r, "  |  ", 10, False, CARD_BORDER)
    rn(p_r, section_name.split()[-1], 10, True, GRAY400)

# ─── PARSE BULLET WITH INLINE MARKERS ───
def parse_bullet_text(tf, items, accent, is_first_p=True):
    for item in items:
        p = tf.paragraphs[0] if is_first_p else tf.add_paragraph()
        is_first_p = False
        p.space_after = Pt(7); p.space_before = Pt(3)
        rn(p, "▪ ", 11, True, accent)
        
        if "**" in item:
            parts = item.split("**")
            for pi, part in enumerate(parts):
                if pi % 2 == 1:
                    rn(p, part, 11, True, WHITE)
                else:
                    rn(p, part, 11, color=GRAY300)
        else:
            rn(p, item, 11, color=GRAY300)

# ─── FLOATING FOOTER ───
def add_footer(slide, active_page, accent):
    rect(slide, 0.8, 6.95, 11.733, 0.015, CARD_BORDER)
    lt = tx(slide, 0.8, 7.02, 6.0, 0.3)
    lp = lt.text_frame.paragraphs[0]
    rn(lp, "MIN DO KIM  ·  BACKEND DEVELOPER PORTFOLIO", 8, True, GRAY500)
    
    pt = tx(slide, 9.0, 7.02, 3.533, 0.3)
    pp = pt.text_frame.paragraphs[0]; pp.alignment = PP_ALIGN.RIGHT
    rn(pp, active_page, 9.5, True, accent)

# ─── SECTION DIVIDER (간지) SLIDE ───
def divider(part_num, title, subtitle, accent, accent_lite, pills_list=None):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, DARK)
    
    rect(sl, 4.6, 0, 0.08, 7.5, accent)
    
    pt = tx(sl, 0.6, 1.8, 3.5, 0.4)
    pp = pt.text_frame.paragraphs[0]
    rn(pp, f"PART 0{part_num}.", 12, True, accent)
    
    nt = tx(sl, 0.6, 2.3, 3.5, 2.5)
    np = nt.text_frame.paragraphs[0]
    rn(np, f"0{part_num}", 100, True, WHITE)
    
    cat_lbls = ["CORE ENTERPRISE", "FOUNDATIONAL PROJECTS"]
    ct = tx(sl, 0.6, 5.0, 3.5, 0.4)
    cp = ct.text_frame.paragraphs[0]
    rn(cp, cat_lbls[part_num - 1], 10.5, True, GRAY400)
    
    tt = tx(sl, 5.2, 2.0, 7.5, 1.0)
    tp = tt.text_frame.paragraphs[0]
    rn(tp, title, 36, True, WHITE)
    
    st = tx(sl, 5.2, 3.0, 7.2, 1.0)
    st.text_frame.word_wrap = True
    sp = st.text_frame.paragraphs[0]
    rn(sp, subtitle, 14, color=GRAY300)
    
    rect(sl, 5.2, 4.3, 7.2, 0.015, CARD_BORDER)
    
    if pills_list:
        draw_tech_pills(sl, 5.2, 4.7, pills_list, None, accent)

# ─── OVERVIEW LAYOUT ───
def overview_slide(project_idx, section, title, items, accent, img_path=None, img_caption=None, active_page=""):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, DARK)
    add_clean_header(sl, project_idx, section, title, accent)
    add_footer(sl, active_page, accent)
    
    specs = []
    desc_items = []
    for item in items:
        if any(x in item for x in ["**개발 기간**", "**개발 인원**", "**담당 역할**", "**사용 기술**", "**프로젝트 개요**", "**기획 의도**"]):
            specs.append(item)
        else:
            desc_items.append(item)
            
    rrect(sl, 0.8, 1.6, 3.8, 5.0, CARD_BG, CARD_BORDER)
    rect(sl, 0.8, 1.6, 3.8, 0.08, accent)
    
    tb_spec = tx(sl, 1.05, 1.85, 3.3, 4.5)
    tf_s = tb_spec.text_frame; tf_s.word_wrap = True
    tf_s.margin_left = tf_s.margin_top = tf_s.margin_right = tf_s.margin_bottom = 0
    
    p_h = tf_s.paragraphs[0]; p_h.space_after = Pt(12)
    rn(p_h, "PROJECT SPEC", 13, True, WHITE)
    
    is_first = True
    for spec in specs:
        p = tf_s.paragraphs[0] if is_first else tf_s.add_paragraph()
        is_first = False
        p.space_after = Pt(8); p.space_before = Pt(2)
        
        if "  " in spec:
            key, val = spec.split("  ", 1)
        elif "** " in spec:
            key, val = spec.split("** ", 1)
            key = key + "**"
        else:
            key, val = spec, ""
            
        k_clean = key.replace("**", "").strip()
        rn(p, f"▪ {k_clean}\n", 10, True, accent)
        rn(p, val.strip(), 9.5, color=GRAY300)
        
    has_img = img_path and os.path.exists(img_path)
    if has_img:
        rrect(sl, 4.9, 1.6, 7.633, 3.2, CARD_BG, CARD_BORDER)
        pic(sl, img_path, 5.05, 1.75, 7.333, 2.5)
        if img_caption:
            tb_cap = tx(sl, 5.05, 4.3, 7.333, 0.3)
            p = tb_cap.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            rn(p, img_caption, 8.5, italic=True, color=GRAY400)
            
        tb_desc = tx(sl, 4.9, 4.9, 7.633, 1.7)
        tf_d = tb_desc.text_frame; tf_d.word_wrap = True
        tf_d.margin_left = tf_d.margin_top = tf_d.margin_bottom = tf_d.margin_right = 0
        parse_bullet_text(tf_d, desc_items, accent, is_first_p=True)
    else:
        tb_desc = tx(sl, 4.9, 1.6, 7.633, 5.0)
        tf_d = tb_desc.text_frame; tf_d.word_wrap = True
        tf_d.margin_left = tf_d.margin_top = tf_d.margin_bottom = tf_d.margin_right = 0
        parse_bullet_text(tf_d, desc_items, accent, is_first_p=True)

# ─── CORE FEATURES GRID LAYOUT ───
def features_slide(project_idx, section, title, items, accent, active_page=""):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, DARK)
    add_clean_header(sl, project_idx, section, title, accent)
    add_footer(sl, active_page, accent)
    
    n = len(items)
    if n <= 4:
        cols = n; rows = 1
    elif n <= 6:
        cols = 3; rows = 2
    else:
        cols = 4; rows = 2
        
    for idx, item in enumerate(items):
        r = idx // cols
        c = idx % cols
        
        if cols == 3:
            w_card = 3.68; gap = 0.35; l_start = 0.8
        elif cols == 4:
            w_card = 2.68; gap = 0.32; l_start = 0.8
        else:
            w_card = 5.6; gap = 0.5; l_start = 0.8
            
        x = l_start + c * (w_card + gap)
        y = 1.6 + r * 2.55
        h_card = 2.3
        
        rrect(sl, x, y, w_card, h_card, CARD_BG, CARD_BORDER)
        rect(sl, x, y, w_card, 0.06, accent)
        
        tb = tx(sl, x + 0.15, y + 0.15, w_card - 0.3, h_card - 0.3)
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        text = item
        if text[0].isdigit() and text[1:3] in [". ", ") ", "  "]:
            text = text[3:]
            
        card_num = f"0{idx + 1}"
        p = tf.paragraphs[0]; p.space_after = Pt(4)
        rn(p, f"{card_num}  ", 10.5, True, accent)
        
        if "**" in text:
            parts = text.split("**")
            title_text = parts[1]
            desc_text = "".join(parts[2:]).strip(" —").strip()
            rn(p, title_text, 11, True, WHITE)
            p2 = tf.add_paragraph(); p2.space_before = Pt(2)
            rn(p2, desc_text, 9.5, color=GRAY300)
        else:
            rn(p, text, 10.5, color=GRAY300)

# ─── TECHNICAL DETAIL LAYOUT ───
def technical_slide(project_idx, section, title, items, accent, img_path=None, img_caption=None, active_page=""):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, DARK)
    add_clean_header(sl, project_idx, section, title, accent)
    add_footer(sl, active_page, accent)
    
    has_img = img_path and os.path.exists(img_path)
    if has_img:
        tb_text = tx(sl, 0.8, 1.6, 5.8, 5.0)
        tf = tb_text.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        parse_bullet_text(tf, items, accent, is_first_p=True)
        
        rrect(sl, 6.9, 1.6, 5.633, 5.0, CARD_BG, CARD_BORDER)
        pic(sl, img_path, 7.05, 1.75, 5.333, 4.0)
        if img_caption:
            tb_cap = tx(sl, 7.05, 5.85, 5.333, 0.35)
            p = tb_cap.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            rn(p, img_caption, 9, italic=True, color=GRAY400)
    else:
        if len(items) >= 5:
            mid = (len(items) + 1) // 2
            left_items = items[:mid]
            right_items = items[mid:]
            
            ct_l = tx(sl, 0.8, 1.6, 5.5, 5.0)
            tf_l = ct_l.text_frame; tf_l.word_wrap = True
            tf_l.margin_left = tf_l.margin_top = tf_l.margin_bottom = tf_l.margin_right = 0
            parse_bullet_text(tf_l, left_items, accent, is_first_p=True)
            
            rect(sl, 6.666, 1.8, 0.015, 4.5, CARD_BORDER)
            
            ct_r = tx(sl, 7.0, 1.6, 5.5, 5.0)
            tf_r = ct_r.text_frame; tf_r.word_wrap = True
            tf_r.margin_left = tf_r.margin_top = tf_r.margin_bottom = tf_r.margin_right = 0
            parse_bullet_text(tf_r, right_items, accent, is_first_p=True)
        else:
            ct = tx(sl, 0.8, 1.6, 11.733, 5.0)
            tf = ct.text_frame; tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
            parse_bullet_text(tf, items, accent, is_first_p=True)

# ─── TROUBLESHOOTING LAYOUT ───
def split_trouble_items(items):
    first = items[0]
    if "**문제**" in first or "**N+1" in first or "배경" in first:
        if len(items) > 1:
            for idx in range(1, len(items)):
                if any(kw in items[idx] for kw in ["**해결**", "**1단계**", "**뒤로가기**", "**공동구매 동시**", "**화면 단**", "**적용**", "**K6 테스트**"]):
                    return items[:idx], items[idx:]
            return items[:1], items[1:]
            
    for idx in range(1, len(items)):
        if any(kw in items[idx] for kw in ["**해결**", "**1단계**", "**뒤로가기**", "**공동구매 동시**", "**화면 단**", "**적용**", "**K6 테스트**", "**Soft Delete**", "**수강생 등록**"]):
            return items[:idx], items[idx:]
            
    mid = (len(items) + 1) // 2
    return items[:mid], items[mid:]

def troubleshooting_slide(project_idx, section, title, items, accent, img_path=None, img_caption=None, active_page=""):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, DARK)
    add_clean_header(sl, project_idx, section, title, accent)
    add_footer(sl, active_page, accent)
    
    left_items, right_items = split_trouble_items(items)
    has_img = img_path and os.path.exists(img_path)
    
    if has_img:
        rrect(sl, 0.8, 1.6, 5.8, 2.3, CARD_BG, CARD_BORDER)
        rect(sl, 0.8, 1.6, 0.08, 2.3, accent)
        tb_p = tx(sl, 1.05, 1.75, 5.35, 2.05)
        tf_p = tb_p.text_frame; tf_p.word_wrap = True
        tf_p.margin_left = tf_p.margin_top = tf_p.margin_bottom = tf_p.margin_right = 0
        p_ph = tf_p.paragraphs[0]; p_ph.space_after = Pt(4)
        rn(p_ph, "PROBLEM & CONTEXT", 11, True, WHITE)
        parse_bullet_text(tf_p, left_items, accent, is_first_p=False)
        
        rrect(sl, 0.8, 4.1, 5.8, 2.5, CARD_BG, CARD_BORDER)
        rect(sl, 0.8, 4.1, 0.08, 2.5, accent)
        tb_s = tx(sl, 1.05, 4.25, 5.35, 2.25)
        tf_s = tb_s.text_frame; tf_s.word_wrap = True
        tf_s.margin_left = tf_s.margin_top = tf_s.margin_bottom = tf_s.margin_right = 0
        p_sh = tf_s.paragraphs[0]; p_sh.space_after = Pt(4)
        rn(p_sh, "SOLUTION & RESOLUTION", 11, True, WHITE)
        parse_bullet_text(tf_s, right_items, accent, is_first_p=False)
        
        rrect(sl, 6.9, 1.6, 5.6, 5.0, CARD_BG, CARD_BORDER)
        pic(sl, img_path, 7.05, 1.75, 5.3, 4.0)
        if img_caption:
            tb_cap = tx(sl, 7.05, 5.85, 5.3, 0.35)
            p = tb_cap.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            rn(p, img_caption, 8.5, italic=True, color=GRAY400)
    else:
        rrect(sl, 0.8, 1.6, 5.6, 5.0, CARD_BG, CARD_BORDER)
        rect(sl, 0.8, 1.6, 0.08, 5.0, accent)
        tb_l = tx(sl, 1.05, 1.8, 5.15, 4.6)
        tf_l = tb_l.text_frame; tf_l.word_wrap = True
        tf_l.margin_left = tf_l.margin_top = tf_l.margin_bottom = tf_l.margin_right = 0
        p_lh = tf_l.paragraphs[0]; p_lh.space_after = Pt(10)
        rn(p_lh, "PROBLEM & CONTEXT", 12, True, WHITE)
        parse_bullet_text(tf_l, left_items, accent, is_first_p=False)
        
        rrect(sl, 6.9, 1.6, 5.6, 5.0, CARD_BG, CARD_BORDER)
        rect(sl, 6.9, 1.6, 0.08, 5.0, accent)
        tb_r = tx(sl, 7.15, 1.8, 5.15, 4.6)
        tf_r = tb_r.text_frame; tf_r.word_wrap = True
        tf_r.margin_left = tf_r.margin_top = tf_r.margin_bottom = tf_r.margin_right = 0
        p_rh = tf_r.paragraphs[0]; p_rh.space_after = Pt(10)
        rn(p_rh, "SOLUTION & RESOLUTION", 12, True, WHITE)
        parse_bullet_text(tf_r, right_items, accent, is_first_p=False)

# ─── RETROSPECTIVE LAYOUT ───
RETRO_QUOTES = {
    0: "인프라와 분산 시스템의 기술적 깊이를 이해하고\n트레이드 오프를 고민할 줄 아는\n개발자로 성장하겠습니다.",
    1: "K6 실측 테스트 기반으로\nRedisson 락 + Redis 캐시의 최적 결합 구조를\n도출하며 동시성 제어의 깊이를 경험했습니다.",
    2: "Spring 없이 Servlet 라이프사이클을\n직접 구현하며 HTTP 흐름과 Filter 보안 동작을\n하부 구조부터 체득했습니다.",
    3: "ERD 모델링과 정규화/반정규화\n트레이드오프를 실질적으로 검증하며\nDB 스키마 설계 역량을 확보했습니다."
}

def retrospective_slide(project_idx, section, title, items, accent, active_page=""):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, DARK)
    add_clean_header(sl, project_idx, section, title, accent)
    add_footer(sl, active_page, accent)
    
    rrect(sl, 0.8, 1.6, 4.8, 5.0, RGBColor(10, 15, 30), None)
    rect(sl, 0.8, 1.6, 0.12, 5.0, accent)
    
    tb_q = tx(sl, 1.2, 2.2, 4.0, 4.0)
    tf_q = tb_q.text_frame; tf_q.word_wrap = True
    tf_q.margin_left = tf_q.margin_top = tf_q.margin_bottom = tf_q.margin_right = 0
    p_q1 = tf_q.paragraphs[0]; p_q1.space_after = Pt(8)
    rn(p_q1, "DEVELOPER'S NOTE", 11, True, accent)
    
    p_q2 = tf_q.add_paragraph(); p_q2.space_before = Pt(8); p_q2.space_after = Pt(12)
    quote_text = RETRO_QUOTES.get(project_idx)
    rn(p_q2, f'"{quote_text}"', 14, False, WHITE, italic=True)
    
    lessons = []
    curr_title = ""
    curr_desc = []
    for item in items:
        if item.startswith("**") and item.endswith("**"):
            if curr_title:
                lessons.append((curr_title, " ".join(curr_desc)))
            curr_title = item.replace("**", "").strip()
            curr_desc = []
        else:
            curr_desc.append(item)
    if curr_title:
        lessons.append((curr_title, " ".join(curr_desc)))
        
    num_lessons = len(lessons)
    for idx, (ltitle, ldesc) in enumerate(lessons[:3]):
        if num_lessons >= 3:
            h_card = 1.5; y = 1.6 + idx * 1.73; text_size = 9.5; desc_space = Pt(2)
        elif num_lessons == 2:
            h_card = 2.38; y = 1.6 + idx * 2.62; text_size = 10; desc_space = Pt(4)
        else:
            h_card = 5.0; y = 1.6; text_size = 10.5; desc_space = Pt(6)
            
        rrect(sl, 6.0, y, 6.5, h_card, CARD_BG, CARD_BORDER)
        rect(sl, 6.0, y, 6.5, 0.06, accent)
        
        tb_l = tx(sl, 6.2, y + 0.15, 6.1, h_card - 0.25)
        tf_l = tb_l.text_frame; tf_l.word_wrap = True
        tf_l.margin_left = tf_l.margin_top = tf_l.margin_bottom = tf_l.margin_right = 0
        
        p = tf_l.paragraphs[0]; p.space_after = desc_space
        rn(p, f"0{idx + 1}  {ltitle}", 11.5, True, WHITE)
        p2 = tf_l.add_paragraph()
        rn(p2, ldesc, text_size, color=GRAY300)

# ─── HIGH-COMPRESSION SUPPORTING PROJECT SUMMARY LAYOUT ───
def summary_project_slide(project_idx, name_ko, spec_items, bullet_items, img_path, caption, active_page, accent):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, DARK)
    add_clean_header(sl, 1, f"FOUNDATIONAL PROJECT 0{project_idx}", f"{PROJ_NAMES_KO[project_idx]} 요약", accent)
    add_footer(sl, active_page, accent)
    
    # Left Spec Card
    rrect(sl, 0.8, 1.6, 4.2, 5.0, CARD_BG, CARD_BORDER)
    rect(sl, 0.8, 1.6, 4.2, 0.08, accent)
    
    tb_spec = tx(sl, 1.05, 1.85, 3.7, 4.5)
    tf_s = tb_spec.text_frame; tf_s.word_wrap = True
    tf_s.margin_left = tf_s.margin_top = tf_s.margin_right = tf_s.margin_bottom = 0
    p_h = tf_s.paragraphs[0]; p_h.space_after = Pt(12)
    rn(p_h, "PROJECT SPECIFICATION", 12.5, True, WHITE)
    
    is_first = True
    for spec in spec_items:
        p = tf_s.paragraphs[0] if is_first else tf_s.add_paragraph()
        is_first = False
        p.space_after = Pt(8); p.space_before = Pt(2)
        
        if "  " in spec:
            key, val = spec.split("  ", 1)
        elif "** " in spec:
            key, val = spec.split("** ", 1)
            key = key + "**"
        else:
            key, val = spec, ""
            
        k_clean = key.replace("**", "").strip()
        rn(p, f"▪ {k_clean}\n", 10, True, accent)
        rn(p, val.strip(), 9.5, color=GRAY300)
        
    # Right Content Area (Top Image, Bottom Key Features)
    has_img = img_path and os.path.exists(img_path)
    
    if has_img:
        # Top image card
        rrect(sl, 5.3, 1.6, 7.233, 2.5, CARD_BG, CARD_BORDER)
        pic(sl, img_path, 5.45, 1.7, 6.933, 1.9)
        if caption:
            tb_cap = tx(sl, 5.45, 3.65, 6.933, 0.3)
            p = tb_cap.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            rn(p, caption, 8.5, italic=True, color=GRAY400)
            
        # Bottom Summary Card
        rrect(sl, 5.3, 4.25, 7.233, 2.35, CARD_BG, CARD_BORDER)
        tb_sm = tx(sl, 5.5, 4.4, 6.833, 2.1)
        tf_sm = tb_sm.text_frame; tf_sm.word_wrap = True
        tf_sm.margin_left = tf_sm.margin_top = tf_sm.margin_right = tf_sm.margin_bottom = 0
        p_sh = tf_sm.paragraphs[0]; p_sh.space_after = Pt(6)
        rn(p_sh, "CORE TECHNICAL ACHIEVEMENTS", 11.5, True, WHITE)
        parse_bullet_text(tf_sm, bullet_items, accent, is_first_p=False)
    else:
        # Full height features card
        rrect(sl, 5.3, 1.6, 7.233, 5.0, CARD_BG, CARD_BORDER)
        tb_sm = tx(sl, 5.5, 1.8, 6.833, 4.6)
        tf_sm = tb_sm.text_frame; tf_sm.word_wrap = True
        tf_sm.margin_left = tf_sm.margin_top = tf_sm.margin_right = tf_sm.margin_bottom = 0
        p_sh = tf_sm.paragraphs[0]; p_sh.space_after = Pt(8)
        rn(p_sh, "CORE TECHNICAL ACHIEVEMENTS", 13, True, WHITE)
        parse_bullet_text(tf_sm, bullet_items, accent, is_first_p=False)


# ══════════════════════════════════════════════════════════════════
# SLIDE 1: COVER
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, DARK)

for gx in [1.5, 3.0, 4.5, 6.0, 7.5, 9.0, 10.5, 12.0]:
    rect(sl, gx, 0, 0.01, 7.5, RGBColor(22, 32, 54), None)
for gy in [1.5, 3.0, 4.5, 6.0]:
    rect(sl, 0, gy, 13.333, 0.01, RGBColor(22, 32, 54), None)

rect(sl, 9.2, 0, 0.8, 7.5, P1_MAIN)
rect(sl, 10.1, 0, 0.8, 7.5, P2_MAIN)
rect(sl, 11.0, 0, 0.8, 7.5, P3_MAIN)
rect(sl, 11.9, 0, 0.8, 7.5, P4_MAIN)

rect(sl, 0.8, 1.8, 0.06, 3.8, P1_MAIN)

mt = tx(sl, 1.1, 1.7, 7.8, 4.5)
tf = mt.text_frame; tf.word_wrap = True
p1 = tf.paragraphs[0]; p1.space_after = Pt(12)
rn(p1, "BACKEND ARCHITECTURE // DEEP-DIVE & SUMMARY", 13, True, P1_MAIN)
p2 = tf.add_paragraph(); p2.space_after = Pt(4)
rn(p2, "KIM MIN DO", 52, True, WHITE)
p3 = tf.add_paragraph(); p3.space_after = Pt(24)
rn(p3, "김민도 백엔드 개발자 포트폴리오 (핵심 집중형)", 22, color=GRAY300)
p4 = tf.add_paragraph(); p4.space_after = Pt(14)
rn(p4, '"끊임없는 질문과 집요함이 프로그램의 설계를 탄탄하게 만듭니다."', 14, color=GRAY400, italic=True)
p5 = tf.add_paragraph(); p5.space_after = Pt(30)
rn(p5, "대규모 트래픽 최적화(골드 마켓) 프로젝트의 아키텍처 및 상세 구현에 역량을 집중하여 입체적으로 구성한 포트폴리오입니다.", 11, color=GRAY300)
p6 = tf.add_paragraph()
rn(p6, "010-3105-7821  ·  kmd192kmd@naver.com  ·  github.com/kmd192kmd", 11.5, color=GRAY400)


# ══════════════════════════════════════════════════════════════════
# SLIDE 2: PROFILE & SKILLS
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, DARK)

add_clean_header(sl, 0, "ABOUT ME", "프로필 및 보유 기술 스택 요약", P1_MAIN)

rrect(sl, 0.8, 1.6, 5.6, 5.0, CARD_BG, CARD_BORDER)
rect(sl, 0.8, 1.6, 5.6, 0.08, P1_MAIN)
lc_t = tx(sl, 1.1, 1.85, 5.0, 4.5)
ltf = lc_t.text_frame; ltf.word_wrap = True
p_h = ltf.paragraphs[0]; p_h.space_after = Pt(14)
rn(p_h, "Profile", 18, True, WHITE)

for label, val in [
    ("이름 / 연령", "김민도  /  1996.06.18"),
    ("학력", "전주대학교 컴퓨터공학과 졸업  (4.12 / 4.5)"),
    ("훈련 이수", "쌍용교육센터 Java Full-Stack 양성과정 (130일, 1040시간)"),
    ("자격증", "정보처리기사 (2025.12)  |  1종보통운전면허"),
    ("소방관 이력", "전직 소방관. 현장의 집요함과 책임감을 설계와 트러블슈팅의 핵심 동력으로 삼고 있습니다."),
]:
    pl = ltf.add_paragraph(); pl.space_before = Pt(4)
    rn(pl, f"▸ {label}", 11, True, WHITE)
    pv = ltf.add_paragraph(); pv.space_after = Pt(2)
    rn(pv, val, 10, color=GRAY300)

rrect(sl, 6.9, 1.6, 5.6, 5.0, CARD_BG, CARD_BORDER)
rect(sl, 6.9, 1.6, 5.6, 0.08, P1_MAIN)
rc_t = tx(sl, 7.2, 1.85, 5.0, 4.5)
rtf = rc_t.text_frame; rtf.word_wrap = True
s_h = rtf.paragraphs[0]; s_h.space_after = Pt(14)
rn(s_h, "Technical Skills", 18, True, WHITE)

for cat, itms in [
    ("Back-end", "Java 21 · Spring Boot · Spring Security · JPA · QueryDSL · MyBatis"),
    ("Database & Cache", "Oracle DB · Redis · Redisson 분산 락 · Elasticsearch"),
    ("Infra & DevOps", "Apache Kafka · Docker · Docker Compose · GitHub Actions CI/CD · AWS EC2"),
    ("Front-end & Tools", "React · Thymeleaf · JavaScript · Git/GitHub · Figma"),
]:
    pc = rtf.add_paragraph(); pc.space_before = Pt(6)
    rn(pc, f"▸ {cat}", 11, True, WHITE)
    pi = rtf.add_paragraph(); pi.space_after = Pt(2)
    rn(pi, itms, 10, color=GRAY300)

add_footer(sl, "02 / 19", P1_MAIN)


# ══════════════════════════════════════════════════════════════════
# SLIDE 3: TABLE OF CONTENTS (Highlight Composition Grid)
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, DARK)

add_clean_header(sl, 0, "INDEX", "수행 프로젝트 인덱스 요약 (핵심 집중형)", P1_MAIN)

# Left Column (Wide) - Gold Market (Core Deep Dive)
rrect(sl, 0.8, 1.6, 6.5, 5.0, CARD_BG, CARD_BORDER)
rect(sl, 0.8, 1.6, 6.5, 0.08, P1_MAIN)
tb_c = tx(sl, 1.1, 1.85, 5.9, 4.5)
tf_c = tb_c.text_frame; tf_c.word_wrap = True
tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0
p = tf_c.paragraphs[0]
rn(p, "PART #1  [CORE DEEP DIVE]", 10.5, True, P1_MAIN)
p2 = tf_c.add_paragraph(); p2.space_before = Pt(6); p2.space_after = Pt(2)
rn(p2, "Gold Market (골드 마켓)", 22, True, WHITE)
p3 = tf_c.add_paragraph(); p3.space_after = Pt(14)
rn(p3, "E-Commerce 쇼핑몰 플랫폼  ·  Spring Boot 3.5 기반", 11.5, True, GRAY400)
rect(sl, 1.1, 2.9, 5.9, 0.015, CARD_BORDER)
p4 = tf_c.add_paragraph(); p4.space_before = Pt(14)
rn(p4, "본 포트폴리오의 대표 프로젝트로서, 대용량 트래픽 최적화와 결제 비동기 격리 아키텍처에 집중하여 상세 기술 성취를 심층 분석합니다.\n\n- Elasticsearch 기반 검색 성능 개선 (250ms ➔ 5~10ms)\n- Apache Kafka 기반 결제 비동기 후속 처리 격리\n- Redis 다중 캐싱, 분산 세션 스토어 및 분산 락 동시성 설계\n- GitHub Actions & Docker 컨테이너 무중단 CI/CD 구축", 10.5, False, GRAY300)

# Right Column (Vertical Stack) - Foundational Projects
supporting = [
    ("PART #2", "Jesiyo (제시요)", "실시간 STOMP 입찰 & Redisson 분산 락", P2_MAIN),
    ("PART #3", "DeveryTime (데버리타임)", "Servlet/JSP MVC Model 2 게시판 구조", P3_MAIN),
    ("PART #4", "교육센터 관리 시스템", "Oracle PL/SQL DB 관계 재정의", P4_MAIN),
]

for idx, (part, name, short_desc, clr) in enumerate(supporting):
    y = 1.6 + idx * 1.7
    rrect(sl, 7.6, y, 4.9, 1.6, CARD_BG, CARD_BORDER)
    rect(sl, 7.6, y, 0.1, 1.6, clr)
    
    tb_s = tx(sl, 7.85, y + 0.15, 4.5, 1.3)
    tf_s = tb_s.text_frame; tf_s.word_wrap = True
    tf_s.margin_left = tf_s.margin_top = tf_s.margin_right = tf_s.margin_bottom = 0
    
    ps1 = tf_s.paragraphs[0]
    rn(ps1, f"{part}   ", 9.5, True, clr)
    rn(ps1, name, 13.5, True, WHITE)
    
    ps2 = tf_s.add_paragraph(); ps2.space_before = Pt(4)
    rn(ps2, short_desc, 9.5, False, GRAY300)

add_footer(sl, "03 / 19", P1_MAIN)


# ══════════════════════════════════════════════════════════════════
# PART 1: GOLD MARKET (DEEP DIVE)
# ══════════════════════════════════════════════════════════════════
p1_techs = ["Spring Boot 3.5", "Elasticsearch", "Kafka", "Redis", "Docker", "AWS EC2", "QueryDSL", "React", "Actions CI/CD"]
divider(1, "Gold Market (골드 마켓)",
        "공동구매·핫딜 서비스를 제공하는 Spring Boot 기반 E-Commerce 쇼핑몰 플랫폼\n선착순 트래픽 하의 동시성 제어, Elasticsearch 검색 최적화, Kafka 비동기 이벤트 격리에 집중",
        P1_MAIN, P1_LINE, p1_techs)

overview_slide(0, "01  개요", "개요 · 기획의도", [
    "**프로젝트 개요**  공동구매·핫딜 서비스를 제공하는 Spring Boot 기반 E-Commerce 플랫폼",
    "**기획 의도**  선착순 공동구매 대규모 트래픽 하에서의 데이터 정합성 보장, Elasticsearch/Redis 기반 조회 최적화, Kafka 비동기 이벤트 격리에 집중",
    "**개발 기간**  2026.05.19 ~ 2026.06.17 (약 4주)",
    "**개발 인원**  4명 (백엔드 메인 프로젝트)",
    "**담당 역할**  상품 검색/조회 · 주문/반품/배송 서비스 · Elasticsearch 오프로딩 · GitHub Actions CI/CD 파이프라인 구축",
    "**사용 기술**  Java 21 · Spring Boot 3.5 · Kafka · Elasticsearch · Redis · Oracle Cloud DB · QueryDSL · React · Docker · AWS EC2",
], P1_MAIN, ip(GMD, "goldmarket-main.png"), "골드마켓 메인 화면", "05 / 19")

features_slide(0, "02  핵심 구현", "핵심 구현 사항 — 요약", [
    "1. **Elasticsearch 검색 오프로딩** — DB 조회를 검색 엔진으로 완전 이관 (250ms → 5~10ms)",
    "2. **Redis 다중 활용** — 인기 상품 캐싱(Cache Warming) · 분산 세션 스토어 · Rate Limiter",
    "3. **Kafka 비동기 이벤트 격리** — 결제 완료 후 후속 처리를 별도 트랜잭션으로 분리하여 장애 전파 차단",
    "4. **공동구매 자동화 대기열** — FIFO + Spring Scheduler로 이탈자 발생 시 다음 대기자 자동 승격",
    "5. **공동구매 동시성 제어** — Redis Rate Limiter(앞단) + DB 비관적 락(최종) 이중 보안 구조",
    "6. **GitHub Actions + Docker CI/CD** — AWS EC2 메모리 부족 문제를 외부 빌드 분리로 해결",
    "7. **React 부분 렌더링** — 주문/취소/반품/환불 다차원 상태 UI를 공통 컴포넌트로 재사용",
], P1_MAIN, "06 / 19")

technical_slide(0, "02  핵심 구현", "Elasticsearch 기반 검색 성능 최적화", [
    "**문제**  정렬·필터 조건 복잡화에 따라 Oracle DB 직접 조회 응답시간 250ms 이상으로 악화",
    "**1단계 — 대표 이미지 비정규화**  매 검색 시 이미지 JOIN 부하를 상품 테이블 컬럼 추가로 제거 (250ms)",
    "**2단계 — DB 인덱싱**  자주 쓰이는 필터/정렬 컬럼에 Oracle 인덱스 설정 (70ms)",
    "**3단계 — Elasticsearch 완전 이관**  검색 연산을 Elasticsearch로 오프로딩하여 DB 부하 제로화 (5~10ms)",
    "**적용 기법**  Painless Script 정렬 튜닝 · Nori 한글 형태소 분석기 · CUD 발생 시 인덱스 비동기 동기화",
], P1_MAIN, ip(GMD, "goldmarket-product-list.png"), "Elasticsearch 기반 상품 목록", "07 / 19")

technical_slide(0, "02  핵심 구현", "Redis 다중 활용 — 캐싱 · 세션 · Rate Limiter", [
    "**인기 상품 추천 캐싱**  서버 기동 시 Cache Warming + 1시간 주기 Scheduler 자동 갱신으로 DB 직접 조회 차단",
    "**분산 세션 스토어**  Spring Session Data Redis로 다중 서버 환경에서 로드밸런서가 다른 서버로 연결해도 로그인 유지",
    "**객체 직렬화 충돌 방지**  서버 배포 시 클래스 변수 불일치 직렬화 예외 → serialVersionUID 고정 + 변환 방식 통일로 차단",
    "**Rate Limiter**  공동구매 참여 시 IP·계정별 반복 요청을 StringRedisTemplate fixed-window 카운터로 앞단 즉시 차단",
], P1_MAIN, ip(GMD, "goldmarket-main.png"), "추천 상품 캐싱 및 분산 세션", "08 / 19")

technical_slide(0, "02  핵심 구현", "Kafka 비동기 이벤트 격리 — 결제 완료 처리", [
    "**문제**  결제 승인·주문 생성·알림이 동기 트랜잭션으로 묶여 후속 실패 시 결제 자체가 롤백",
    "**1단계 — 트랜잭션 1차 완료**  결제 데이터 저장 트랜잭션을 먼저 성공 커밋 (결제 확정)",
    "**2단계 — 이벤트 발행**  @TransactionalEventListener로 커밋 성공 감지 → Apache Kafka로 메시지 발행",
    "**3단계 — 비동기 수신**  Consumer가 주문 생성·알림을 독립 트랜잭션으로 처리",
    "**효과**  후속 기능 장애가 결제로 역류하지 않는 비동기 이벤트 격리 달성",
], P1_MAIN, ip(GMD, "goldmarket-kafka.jpg"), "Kafka 결제 비동기 이벤트 분리도", "09 / 19")

technical_slide(0, "03  데이터 구조", "배송 테이블 통합 & 외부 API 호출 최소화", [
    "**주문/발주 배타적 참조 제약조건 (Exclusive OR)**",
    "- 고객 주문(order_seq)과 업체 발주(purchase_seq) 중 한쪽만 필수 연결되는 배타적 제약 설계",
    "- 단일 배송 테이블로 양 도메인 통합 관리 → 배송 상태 추적 쿼리 단순화",
    "**외부 API 호출 최소화**",
    "- 위도/경도: 주소 등록 시 카카오맵 API 최초 1회 조회 후 DB 영구 저장 → 이후 내부 연산만 수행",
    "- 공휴일 테이블: 한 달 1회 스케줄러로 공공 API 동기화 → 매 쿼리마다 외부 통신 제거",
], P1_MAIN, ip(GMD, "goldmarket-delivery-erd.png"), "배송 ERD 모델링", "10 / 19")

troubleshooting_slide(0, "04  트러블슈팅", "React 부분 렌더링 & 뒤로가기 캐시 문제", [
    "**React 도입 배경**  주문 목록·배송 추적·취소/반품/교환/환불 내역에서 동일 UI 패턴 반복 → 공통 컴포넌트화로 재사용",
    "- 다차원 주문 상태에 따른 분기 렌더링을 State/Props로 처리, JSON API 비동기 통신 연동",
    "**뒤로가기 캐시 미반영 문제** (React 미적용 Thymeleaf 일반 페이지)",
    "- 브라우저 뒤로가기 기본 동작(캐시 복원)을 History API로 가로채어 정지",
    "- pageshow 이벤트 감지 → 최근 본 상품 영역만 Fetch API로 비동기 재호출 → 즉각 데이터 일관성 달성",
], P1_MAIN, ip(GMD, "goldmarket-orders-delivery.png"), "React 기반 주문/배송 추적 UI", "11 / 19")

troubleshooting_slide(0, "04  트러블슈팅", "공동구매 자동 대기열 & 동시성 제어", [
    "**공동구매 자동 참여 전환 대기열**",
    "- 선착순 참여자 이탈 시 FIFO 대기열 첫 번째 대기자에게 24시간 결제 기회 자동 부여",
    "- 기한 초과 시 자동 회수 → 다음 대기자에게 이전 (Spring Scheduler + JPA 트랜잭션 전파)",
    "**공동구매 동시 참여 이중 보안**",
    "- 앞단 Rate Limiter: Redis StringRedisTemplate fixed-window 카운터로 IP·계정별 반복 즉시 차단",
    "- 최종 진입: DB Pessimistic Lock으로 동시 트랜잭션 순차 처리 → 재고 초과 판매 원천 방지",
], P1_MAIN, ip(GMD, "goldmarket-groupbuy.png"), "공동구매 FIFO 대기열", "12 / 19")

troubleshooting_slide(0, "04  트러블슈팅", "주문 취소 동시성 & 환불 멱등성 보장", [
    "**문제**  연속 클릭 중복 환불 위험 + Spring Batch 배송 배치와 취소 요청 충돌",
    "**화면 단**  취소 버튼 즉시 disabled 처리로 물리적 중복 요청 차단",
    "**DB 비관적 락**  백엔드 진입 즉시 Pessimistic Lock 획득 → 충돌 요청 순차 처리",
    "**영속성 컨텍스트 최신화**  Spring Batch 반영 최신 상태를 EntityManager.refresh()로 강제 동기화 후 취소 가능 여부 검증",
    "**외부 결제 API 멱등성**  이미 취소된 결제 건(ALREADY_CANCELED_PAYMENT)은 성공으로 흡수 리턴",
], P1_MAIN, ip(GMD, "goldmarket-cancel.png"), "주문 취소 완료 화면", "13 / 19")

troubleshooting_slide(0, "04  트러블슈팅", "GitHub Actions CI/CD & AWS EC2 배포", [
    "**문제**  RAM 8GiB AWS EC2에서 직접 Gradle 빌드 시 메모리 부족으로 서비스 중단",
    "**해결**  빌드를 GitHub Actions Runner로 외부 분리",
    "- Actions에서 JAR 빌드 → Docker 이미지 생성 → Docker Hub Push",
    "- EC2에서 이미지 Pull → Docker Compose 컨테이너 기동 → 서버 메모리 빌드 부하 제로화",
    "**보안 사고 경험**  Git에 API Key 노출 → git-filter-repo로 History 수정 + .env 환경 변수 격리",
    "**자동화 파이프라인**  main Push → Actions 자동 빌드 → Docker Hub → EC2 SSH 접속 → Compose 재시작",
], P1_MAIN, ip(GMD, "goldmarket-cicd.jpg"), "CI/CD 구성 도식", "14 / 19")

troubleshooting_slide(0, "04  트러블슈팅", "매출 통계 N+1 최적화 & 대량 쿠폰 발급", [
    "**N+1 쿼리 병목**  매출 대시보드 조회 시 결제·상품 정보를 개별 요청으로 가져오는 다중 쿼리 발생",
    "**해결**  QueryDSL Dynamic 조건문을 단일 JOIN 쿼리로 통합 → 조회 속도 개선",
    "**BigDecimal 정밀 연산**  금융 정산의 소수점 오차 예방을 위해 double 대신 BigDecimal 사용",
    "**대량 쿠폰 발급**  수천 건 단건 INSERT 루프 → JDBC Batch Insert로 일괄 처리하여 속도 획기적 향상",
], P1_MAIN, ip(GMD, "goldmarket-admin-stats.png"), "관리자 매출 통계 대시보드", "15 / 19")

retrospective_slide(0, "05  회고", "프로젝트 완수 후기 · 회고", [
    "**배포 및 보안 경험**",
    "GitHub Actions CI/CD 구축 중 메모리 부족 문제와 AWS EC2 배포 오류들을 하나씩 해결하며 서버·배포 실전 경험을 쌓았습니다. 특히 Git에 API Key 노출 사고를 겪으며 Git History 수정 및 환경 변수 격리의 중요성을 직접 체험했습니다.",
    "**바이브 코딩 시대와 트레이드 오프 설계**",
    "AI가 코드를 빠르게 완성해 주는 시대지만, 시스템 구조 설계·검색 성능 최적화 시 어느 수준까지 고도화할지, 비용 대비 효과를 결정하는 트레이드 오프 판단은 결국 개발자의 몫입니다. 인프라와 분산 시스템의 기술적 깊이를 이해하고 트레이드 오프를 고민할 줄 아는 개발자로 성장하겠습니다.",
], P1_MAIN, "16 / 19")


# ══════════════════════════════════════════════════════════════════
# PART 2: FOUNDATIONAL PROJECTS (HIGH COMPRESSION SUMMARY)
# ══════════════════════════════════════════════════════════════════
divider(2, "Foundational Projects",
        "Jesiyo, DeveryTime, 교육센터 3대 기초 역량 프로젝트 요약\n기본적인 백엔드 웹 아키텍처, 데이터 무결성 DDL 설계 및 동시성 제어 기본기 검증",
        P2_MAIN, P2_LINE, ["Java 11/21", "Servlet/JSP", "Oracle Database", "WebSocket", "Redisson Lock"])

PROJ_NAMES_KO = ["Gold Market (골드 마켓)", "Jesiyo (제시요)", "DeveryTime (데버리타임)", "교육센터 운영 관리 시스템"]

# Jesiyo Summary Slide
summary_project_slide(
    project_idx=1,
    name_ko=PROJ_NAMES_KO[0],
    spec_items=[
        "**개발 기간**  2026.04.15 ~ 2026.04.29 (15일)",
        "**개발 인원**  4명",
        "**담당 역할**  경매·실시간 입찰 도메인 전담, Redisson 분산 락, Redis 캐싱, K6 테스트",
        "**사용 기술**  Oracle 19c · Java 11 · Spring 5.3 · Redis · Redisson · K6 · MyBatis"
    ],
    bullet_items=[
        "**WebSocket/STOMP 실시간 입찰·채팅** — 채팅방 내부 경매 탭 통합 UI 구현으로 UX 오버헤드 제거",
        "**Redisson 분산 락 동시성 제어** — 대량 입찰 트래픽 하에서의 최고가 정합성 및 무결성 보장",
        "**Redis 최고가 캐시 (Cache-Aside)** — K6 부하 테스트를 거쳐 최고가 조회를 Redis로 오프로딩 (응답률 98.61% 확보, TPS 9배 향상)"
    ],
    img_path=ip(EXT, "image31.png"),
    caption="제시요 실시간 경매 및 채팅 연동 화면",
    active_page="18 / 19",
    accent=P2_MAIN
)

# DeveryTime Summary Slide
summary_project_slide(
    project_idx=2,
    name_ko=PROJ_NAMES_KO[1],
    spec_items=[
        "**개발 기간**  2026.03.19 ~ 2026.03.31 (13일)",
        "**개발 인원**  4명",
        "**담당 역할**  스터디 도메인 관리, 권한/인가 Filter, To-Do AJAX, FullCalendar.js",
        "**사용 기술**  Oracle 21c · Tomcat 9 · Servlet/JSP (Model 2) · JDBC · FullCalendar.js"
    ],
    bullet_items=[
        "**JSP Model 2 MVC 아키텍처** — 프레임워크 없이 Controller-Service-DAO 계층을 바닥부터 구현",
        "**Fetch API 비동기 To-Do** — 서버 응답 성공 시에만 체크 완료 반영 및 통신 실패 시 상태 원상 복구(Rollback)",
        "**데이터베이스 무결성 설계** — 복합 Unique 제약으로 중복 방지, soft delete 대응 PL/SQL 트리거 구현"
    ],
    img_path=ip(EXT, "image20.png"),
    caption="데버리타임 메인 UI 화면",
    active_page="19 / 19",
    accent=P3_MAIN
)

# 교육센터 Summary Slide
summary_project_slide(
    project_idx=3,
    name_ko=PROJ_NAMES_KO[2],
    spec_items=[
        "**개발 기간**  2026.02.03 ~ 2026.02.10 (8일)",
        "**개발 인원**  4명",
        "**담당 역할**  전체 테이블 및 제약조건 설계, 의존성 트리, 정원 초과 방지 트리거",
        "**사용 기술**  Oracle 21c · PL/SQL · ERDCloud · Figma"
    ],
    bullet_items=[
        "**5단계 테이블 의존성 트리** — 외래키(FK) 참조 방향 분석을 통한 대량 더미 데이터 순차 삽입",
        "**정규화 vs 반정규화 고찰** — 시험 배점 필기/실기 컬럼 분리 모델링으로 조회 쿼리 단순화",
        "**PL/SQL 트리거 & CHECK 제약** — 배점 합 CHECK 제약조건 도입 및 강의실 정원 초과 방지 트리거 구현"
    ],
    img_path=ip(EXT, "image8.png"),
    caption="테이블 의존성 트리 설계 도식",
    active_page="19 / 19", # (Divider가 추가되어 19 / 19로 표시되도록 조정)
    accent=P4_MAIN
)


# ══════════════════════════════════════════════════════════════════
# CLOSING SLIDE
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, DARK)

rect(sl, 9.2, 0, 0.8, 7.5, P1_MAIN)
rect(sl, 10.1, 0, 0.8, 7.5, P2_MAIN)
rect(sl, 11.0, 0, 0.8, 7.5, P3_MAIN)
rect(sl, 11.9, 0, 0.8, 7.5, P4_MAIN)

rect(sl, 0.8, 2.2, 0.06, 3.2, P1_MAIN)

mt = tx(sl, 1.1, 2.0, 7.8, 4.5)
tf = mt.text_frame; tf.word_wrap = True

p1 = tf.paragraphs[0]; p1.space_after = Pt(14)
rn(p1, "THANK YOU.", 14, True, P1_MAIN)
p2 = tf.add_paragraph(); p2.space_after = Pt(24)
rn(p2, "감사합니다.", 50, True, WHITE)
p3 = tf.add_paragraph(); p3.space_after = Pt(30)
rn(p3, '"현장의 안전을 책임지던 소방관의 집요함으로\n기초 모델링부터 동시성 제어까지 단계별로 성장해 왔습니다."', 17, color=GRAY300, italic=True)
p4 = tf.add_paragraph()
rn(p4, "010-3105-7821  ·  kmd192kmd@naver.com  ·  github.com/kmd192kmd", 12.5, color=GRAY500)


# ══════════════════════════════════════════════════════════════════
# SAVE PRESENTATION
# ══════════════════════════════════════════════════════════════════
out_path = r"C:\Users\a\Desktop\쇼핑몰 프로젝트 자료모음\6_자소서, 포폴 1차 완성본_260525\포트폴리오(김민도)_v10.pptx"
prs.save(out_path)
print(f"SUCCESS: Saved portfolio to {out_path}")
print(f"Slide count: {len(prs.slides)}")
